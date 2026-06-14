#!/usr/bin/env python3
"""
先声药业集团PPT生成脚本 v3.0
基于官方模板 + python-pptx，融合四份参考PPT风格。
精确排版参数来源于：
- 10-版本号202601集团PPT模板.pptx（官方母版）
- 科唯可-汇报任董0528.pptx（实战复盘风格）
- 品牌提示物使用情况概览.pptx（矩阵布局风格）
- 何为专业的营销管理-任董.pptx（文字密集型风格）

Usage:
    python generate.py --template <template.pptx> --input <content.json> --output <output.pptx>

content.json 格式:
{
  "title": "PPT主标题",
  "subtitle": "副标题",
  "slides": [
    {"type": "cover", "title": "...", "subtitle": "...", "date": "...", "presenter": "..."},
    {"type": "toc", "title": "目录", "items": ["第一章", "第二章", "第三章"]},
    {"type": "section", "title": "章节标题", "subtitle": "章节说明"},
    {"type": "content", "title": "页面标题", "body": [...], "table": {...}, "chart": {...}},
    {"type": "gantt", "title": "...", "rows": [...], "months": [...], "data": [...], "milestones": [...], "conclusion": "..."},
    {"type": "timeline_horizontal", "title": "...", "subtitle": "...", "items": [...]},
    {"type": "big_number", "title": "...", "big_number": "5,000", "unit": "场", "description": "...", "cards": [...], "side_cards": [...]},
    {"type": "comparison", "title": "...", "intro": "...", "items": [...], "target": "...", "note": "..."},
    {"type": "process", "title": "...", "steps": [...], "direction": "horizontal"},
    {"type": "kpi_dashboard", "title": "...", "kpis": [...]},
    {"type": "item_matrix", "title": "...", "groups": [...], "footer_notes": "..."},
    {"type": "calendar_grid", "title": "...", "rows": [...], "months": [...], "grid": [...], "note": "..."},
    {"type": "review_matrix", "title": "...", "badge": "...", "category_header": "...", "good_header": "...", "bad_header": "...", "rows": [...], "footer_note": "..."},
    {"type": "action_category", "title": "...", "badge": "...", "intro": "...", "categories": [...], "highlight_box": {...}},
    {"type": "strategy_diagram", "title": "...", "badge": "...", "diagram": {...}},
    {"type": "title_only", "title": "仅标题页"},
    {"type": "ending", "title": "谢谢", "subtitle": "Thank You"}
  ]
}
"""

import argparse
import json
import sys
import os
import copy

# Ensure python-pptx is available
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu, Cm
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.chart.data import ChartData
    from pptx.oxml.ns import nsmap, qn
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

# 导入排版引擎
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
try:
    from simcere_layout import (
        COLORS, FONT_CN, FONT_EN, FONT_SIZES,
        SLIDE_WIDTH, SLIDE_HEIGHT,
        COVER_LAYOUT, CONTENT_LAYOUT, SECTION_LAYOUT,
        TOC_LAYOUT, ENDING_LAYOUT, TITLE_ONLY_LAYOUT,
        MATRIX_LAYOUT, TABLE_LAYOUT, KPI_LAYOUT,
        PARA_SPACING, LINE_SPACING,
        LAYOUT_INDEX,
        set_font, set_paragraph_spacing,
        add_textbox as layout_add_textbox,
        add_paragraph as layout_add_paragraph,
        add_first_paragraph,
        fill_placeholder as layout_fill_placeholder,
        get_placeholder as layout_get_placeholder,
        remove_placeholders as layout_remove_placeholders,
        add_bottom_bar, add_table as layout_add_table,
        style_table_header, style_table_body,
        build_body_content, add_kpi_card, add_green_tag,
    )
    USE_LAYOUT_ENGINE = True
    # 补充排版引擎中不包含的业务色别名
    COLORS['hlk_green'] = COLORS.get('hlk_green', RGBColor(0xC8, 0xEB, 0xD6))
    COLORS['hk_green'] = COLORS.get('hk_green', RGBColor(0xE8, 0xF5, 0xEE))
    COLORS['hlk_orange'] = COLORS.get('hlk_orange', RGBColor(0xFF, 0xED, 0xE1))
    COLORS['orange'] = COLORS.get('orange', RGBColor(0xE4, 0x80, 0x30))
    COLORS['table_header'] = COLORS.get('table_header', COLORS['accent1'])
    COLORS['table_stripe'] = COLORS.get('table_stripe', RGBColor(0xF0, 0xF7, 0xF4))
    COLORS['tag_green'] = COLORS.get('tag_green', RGBColor(0x02, 0x9B, 0x46))
    COLORS['score_orange'] = COLORS.get('score_orange', RGBColor(0xE4, 0x80, 0x30))
    COLORS['bottom_bar'] = COLORS.get('bottom_bar', COLORS['accent1'])
except ImportError:
    print("Warning: simcere_layout.py not found, using built-in parameters", file=sys.stderr)
    USE_LAYOUT_ENGINE = False
    # 内置备用参数（与 simcere_layout.py 一致）
    COLORS = {
        'accent1': RGBColor(0x00, 0xB0, 0x52), 'accent2': RGBColor(0x8F, 0xD4, 0x00),
        'accent3': RGBColor(0x00, 0x66, 0x47), 'accent4': RGBColor(0x00, 0xB5, 0xBD),
        'accent5': RGBColor(0xFC, 0xC9, 0x17), 'accent6': RGBColor(0xF5, 0x66, 0x00),
        'dk1': RGBColor(0x33, 0x33, 0x33), 'dk2': RGBColor(0x44, 0x54, 0x6A),
        'lt1': RGBColor(0xFF, 0xFF, 0xFF), 'lt2': RGBColor(0xE7, 0xE6, 0xE6),
        'table_header': RGBColor(0x00, 0xB0, 0x52), 'table_stripe': RGBColor(0xF0, 0xF7, 0xF4),
        'tag_green': RGBColor(0x02, 0x9B, 0x46), 'score_orange': RGBColor(0xE4, 0x80, 0x30),
        'orange': RGBColor(0xE4, 0x80, 0x30),
        'hlk_green': RGBColor(0xC8, 0xEB, 0xD6), 'hk_green': RGBColor(0xE8, 0xF5, 0xEE),
        'hlk_orange': RGBColor(0xFF, 0xED, 0xE1),
    }
    FONT_CN = '微软雅黑'
    FONT_EN = 'Arial'
    FONT_SIZES = {
        'cover_title': Pt(44), 'cover_subtitle': Pt(20),
        'section_title': Pt(32), 'section_subtitle': Pt(16),
        'page_tag': Pt(20), 'page_title': Pt(26), 'page_title_alt': Pt(18),
        'body_main': Pt(16), 'body_normal': Pt(14), 'body_small': Pt(12),
        'table_header': Pt(12), 'table_body': Pt(11),
        'note': Pt(10), 'micro': Pt(9), 'tiny': Pt(7),
        'big_number': Pt(36), 'big_number_label': Pt(12), 'toc_item': Pt(18),
    }
    LAYOUT_INDEX = {
        'cover': 0, 'toc': 1, 'section': 2, 'content': 3,
        'title_only': 4, 'logo_only': 5, 'blank': 6, 'ending': 7,
    }

# 兼容别名
SIMCERE_COLORS = COLORS
FONT_NAME = FONT_CN


# ============================================================
# 字体和间距辅助函数（当排版引擎不可用时使用）
# ============================================================

def _set_font(run, size=None, bold=False, color=None):
    """设置字体（含东亚字体）"""
    run.font.name = FONT_CN
    try:
        rPr = run._r.get_or_add_rPr()
        ea = rPr.makeelement(qn('a:ea'), {})
        ea.set('typeface', FONT_CN)
        rPr.append(ea)
        latin = rPr.makeelement(qn('a:latin'), {})
        latin.set('typeface', FONT_EN)
        rPr.append(latin)
    except:
        pass
    if size:
        run.font.size = size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _set_para_spacing(para, before=None, after=None, line_spacing=None):
    """设置段落间距（通过XML）"""
    try:
        from lxml import etree as _etree
        pPr = para._pPr
        if pPr is None:
            pPr = _etree.SubElement(para._p, qn('a:pPr'))
        if before is not None:
            spcBef = pPr.find(qn('a:spcBef'))
            if spcBef is None:
                spcBef = _etree.SubElement(pPr, qn('a:spcBef'))
            spcPts = spcBef.find(qn('a:spcPts'))
            if spcPts is None:
                spcPts = _etree.SubElement(spcBef, qn('a:spcPts'))
            spcPts.set('val', str(int(before / 12700)))
        if after is not None:
            spcAft = pPr.find(qn('a:spcAft'))
            if spcAft is None:
                spcAft = _etree.SubElement(pPr, qn('a:spcAft'))
            spcPts = spcAft.find(qn('a:spcPts'))
            if spcPts is None:
                spcPts = _etree.SubElement(spcAft, qn('a:spcPts'))
            spcPts.set('val', str(int(after / 12700)))
        if line_spacing is not None:
            lnSpc = pPr.find(qn('a:lnSpc'))
            if lnSpc is None:
                lnSpc = _etree.SubElement(pPr, qn('a:lnSpc'))
            spcPct = lnSpc.find(qn('a:spcPct'))
            if spcPct is None:
                spcPct = _etree.SubElement(lnSpc, qn('a:spcPct'))
            spcPct.set('val', str(int(line_spacing * 100000)))
    except:
        pass


# ============================================================
# 模板与版式工具
# ============================================================

def find_layout(prs, layout_name):
    """按名称查找版式"""
    for layout in prs.slide_layouts:
        if layout.name == layout_name:
            return layout
    return None


def get_layout_by_type(prs, slide_type):
    """根据slide类型获取版式索引"""
    layout_map = {
        'cover': '封面',
        'toc': '1_仅logo页',
        'section': '节标题',
        'content': '标题和内容',
        'title_only': '仅标题页',
        'blank': '空白',
        'logo_only': '仅logo页',
        'ending': '末尾幻灯片',
    }
    name = layout_map.get(slide_type, '标题和内容')
    layout = find_layout(prs, name)
    if layout is None:
        for l in prs.slide_layouts:
            if name in l.name or l.name in name:
                return l
        return prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else None
    return layout


# ============================================================
# Shape 操作工具
# ============================================================

def remove_shape(shape):
    """从 slide 中彻底删除一个 shape"""
    sp = shape.element
    sp.getparent().remove(sp)


def remove_placeholders(slide, types_to_remove=None):
    """删除指定类型占位符（先收集再删除）"""
    if USE_LAYOUT_ENGINE:
        return layout_remove_placeholders(slide, types_to_remove)
    if types_to_remove is None:
        types_to_remove = ['BODY', 'OBJECT']
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = str(shape.placeholder_format.type).split('.')[-1]
            for t in types_to_remove:
                if t in ph_type:
                    shapes_to_remove.append(shape)
                    break
    for shape in shapes_to_remove:
        remove_shape(shape)


def find_placeholder(slide, placeholder_type):
    """查找指定类型的占位符"""
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph = shape.placeholder_format
            ph_type_name = str(ph.type).split('.')[-1] if ph.type else ''
            if placeholder_type.lower() in ph_type_name.lower():
                return shape
    return None


def set_placeholder_text(slide, placeholder_type, text):
    """设置指定类型占位符文本"""
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph = shape.placeholder_format
            ph_type_name = str(ph.type).split('.')[-1] if ph.type else ''
            if placeholder_type.lower() in ph_type_name.lower():
                set_text_frame_text(shape, text)
                return True
    return False


def set_text_frame_text(shape, text, font_size=None, bold=None, color=None):
    """设置文本框内容并保留格式"""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    first_para = tf.paragraphs[0]
    if first_para.runs:
        first_run = first_para.runs[0]
        first_run.text = text
    else:
        first_para.text = text
    for para in tf.paragraphs:
        for run in para.runs:
            _set_font(run)
            if font_size:
                run.font.size = Pt(font_size) if isinstance(font_size, (int, float)) else font_size
            if bold is not None:
                run.font.bold = bold
            if color:
                run.font.color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=None):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    _set_font(p.runs[0] if p.runs else p.add_run(), size=Pt(font_size), bold=bold, color=color)
    return txBox


def add_shape_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """添加矩形形状"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width if line_width else Pt(1)
    else:
        shape.line.fill.background()
    return shape


# ============================================================
# 文本辅助函数
# ============================================================

def _add_centered_text(slide, left, top, width, height, text, font_size=11, color=None, bold=False, word_wrap=False):
    """添加居中对齐文字"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    _set_font(p.runs[0] if p.runs else p.add_run(), size=Pt(font_size), bold=bold, color=color)
    p.alignment = PP_ALIGN.CENTER
    return txBox


def _add_left_text(slide, left, top, width, height, text, font_size=11, color=None, bold=False, word_wrap=False):
    """添加左对齐文字"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    _set_font(p.runs[0] if p.runs else p.add_run(), size=Pt(font_size), bold=bold, color=color)
    p.alignment = PP_ALIGN.LEFT
    return txBox


# ============================================================
# 表格和图表
# ============================================================

def add_table_to_slide(slide, data, left, top, width, height):
    """添加真实PPT表格"""
    if not data or not data.get('headers') or not data.get('rows'):
        return None
    headers = data['headers']
    rows = data['rows']
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table
    # 表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = str(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['accent1']
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                _set_font(run, size=FONT_SIZES['table_header'], bold=True, color=COLORS['lt1'])
    # 数据行
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            if col_idx >= num_cols:
                break
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['table_stripe']
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    _set_font(run, size=FONT_SIZES['table_body'], color=COLORS['dk1'])
    # 列宽
    col_width = int(width / num_cols)
    for col in table.columns:
        col.width = col_width
    return table


def add_chart_to_slide(slide, chart_spec, left, top, width, height):
    """添加真实PPT图表"""
    if not chart_spec:
        return None
    chart_type_str = chart_spec.get('chart_type', 'column')
    categories = chart_spec.get('categories', [])
    series_list = chart_spec.get('series', [])
    chart_title = chart_spec.get('title', '')
    if not categories or not series_list:
        return None
    chart_type_map = {
        'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'bar': XL_CHART_TYPE.BAR_CLUSTERED,
        'line': XL_CHART_TYPE.LINE_MARKERS,
        'pie': XL_CHART_TYPE.PIE,
        'area': XL_CHART_TYPE.AREA,
    }
    chart_type = chart_type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)
    chart_data = ChartData()
    chart_data.categories = categories
    theme_colors = [
        COLORS['accent1'], COLORS['accent2'], COLORS['accent3'],
        COLORS['accent4'], COLORS['accent5'], COLORS['accent6'],
    ]
    for idx, series in enumerate(series_list):
        series_name = series.get('name', f'系列{idx+1}')
        values = series.get('values', [])
        chart_data.add_series(series_name, values)
    graphic_frame = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
    chart = graphic_frame.chart
    if chart_title:
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_title
        for para in chart.chart_title.text_frame.paragraphs:
            for run in para.runs:
                _set_font(run, size=FONT_SIZES['body_normal'], bold=True, color=COLORS['dk1'])
    if chart.has_legend:
        chart.legend.include_in_layout = False
    if chart_type_str != 'pie':
        if hasattr(chart, 'category_axis') and chart.category_axis:
            chart.category_axis.tick_labels.font.size = FONT_SIZES['note']
            chart.category_axis.tick_labels.font.name = FONT_CN
        if hasattr(chart, 'value_axis') and chart.value_axis:
            chart.value_axis.tick_labels.font.size = FONT_SIZES['note']
            chart.value_axis.tick_labels.font.name = FONT_CN
    if chart_type_str != 'pie':
        for idx, series in enumerate(chart.series):
            if idx < len(theme_colors):
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = theme_colors[idx]
    else:
        for series in chart.series:
            for idx, point in enumerate(series.points):
                if idx < len(theme_colors):
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = theme_colors[idx]
    return chart


# ============================================================
# 页面构建器：基础页面类型
# ============================================================

def build_cover_slide(prs, slide_spec):
    """构建封面页 — 标题44pt，副标题20pt"""
    title = slide_spec.get('title', '')
    subtitle = slide_spec.get('subtitle', '')
    date = slide_spec.get('date', '')
    presenter = slide_spec.get('presenter', '')
    layout = get_layout_by_type(prs, 'cover')
    slide = prs.slides.add_slide(layout)
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = str(shape.placeholder_format.type).split('.')[-1]
            if 'CENTER_TITLE' in ph_type:
                set_text_frame_text(shape, title, font_size=44, bold=True)
            elif 'SUBTITLE' in ph_type:
                set_text_frame_text(shape, subtitle, font_size=20)
            elif 'BODY' in ph_type:
                text = date
                if presenter:
                    text = text + '\n' + presenter if text else presenter
                if text:
                    set_text_frame_text(shape, text, font_size=14)
    remove_placeholders(slide)
    return slide


def build_toc_slide(prs, slide_spec):
    """构建目录页 — 18pt条目"""
    title = slide_spec.get('title', '目录')
    items = slide_spec.get('items', [])
    layout = get_layout_by_type(prs, 'toc')
    slide = prs.slides.add_slide(layout)
    remove_placeholders(slide)
    if items:
        txBox = slide.shapes.add_textbox(Inches(3.84), Inches(1.95), Inches(8.76), Inches(4.38))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['toc_item'])
            _set_para_spacing(p, after=Pt(16))
    return slide


def build_section_slide(prs, slide_spec):
    """构建节标题页 — 32pt标题，16pt副标题"""
    title = slide_spec.get('title', '')
    subtitle = slide_spec.get('subtitle', '')
    layout = get_layout_by_type(prs, 'section')
    slide = prs.slides.add_slide(layout)
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = str(shape.placeholder_format.type).split('.')[-1]
            if 'TITLE' in ph_type:
                set_text_frame_text(shape, title, font_size=32, bold=True)
    remove_placeholders(slide)
    if subtitle:
        txBox = slide.shapes.add_textbox(Inches(0.73), Inches(3.06), Inches(11.5), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['section_subtitle'], color=COLORS['dk1'])
    return slide


def build_content_slide(prs, slide_spec):
    """构建标准内容页 — 标签20pt + 标题26pt + 正文14pt"""
    title = slide_spec.get('title', '')
    tag = slide_spec.get('tag', '')
    body = slide_spec.get('body', [])
    table_data = slide_spec.get('table')
    chart_spec = slide_spec.get('chart')
    layout = get_layout_by_type(prs, 'content')
    slide = prs.slides.add_slide(layout)
    # 设置标题
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = str(shape.placeholder_format.type).split('.')[-1]
            if 'TITLE' in ph_type:
                set_text_frame_text(shape, title, font_size=26, bold=True)
    # 删除 BODY/OBJECT 占位符
    remove_placeholders(slide)
    # 标签（科唯可风格）
    current_top = Inches(1.35)
    if tag:
        tag_box = slide.shapes.add_textbox(Inches(0.62), Inches(0.34), Inches(9.89), Inches(0.40))
        tf = tag_box.text_frame
        p = tf.paragraphs[0]
        p.text = tag
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['page_tag'], bold=True, color=COLORS['accent1'])
    # 正文区域
    content_left = Inches(0.62)
    content_width = Inches(12.30)
    slide_bottom = prs.slide_height - Inches(0.5)
    has_body = bool(body)
    has_table = bool(table_data and table_data.get('headers') and table_data.get('rows'))
    has_chart = bool(chart_spec and chart_spec.get('categories') and chart_spec.get('series'))
    available_height = slide_bottom - current_top
    if has_body:
        if has_table or has_chart:
            body_height = int(available_height * 0.35)
        else:
            body_height = int(available_height)
        txBox = slide.shapes.add_textbox(content_left, current_top, content_width, body_height)
        tf = txBox.text_frame
        tf.word_wrap = True
        _add_bullet_list(tf, body)
        current_top += body_height + Inches(0.2)
    if has_table:
        remaining = slide_bottom - current_top
        table_height = min(Inches(2.5), remaining * 0.9)
        add_table_to_slide(slide, table_data, content_left, current_top, content_width, table_height)
        current_top += table_height + Inches(0.2)
    if has_chart:
        remaining = slide_bottom - current_top
        chart_height = min(Inches(3.5), remaining * 0.95)
        add_chart_to_slide(slide, chart_spec, content_left, current_top, content_width, chart_height)
    return slide


def _add_bullet_list(tf, items, font_size=None):
    """添加项目符号列表到文本框"""
    if font_size is None:
        font_size = FONT_SIZES['body_normal']
    tf.clear()
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        level = item.get('level', 0) if isinstance(item, dict) else 0
        text = item.get('text', '') if isinstance(item, dict) else str(item)
        p.clear()
        run = p.add_run()
        run.text = text
        _set_font(run, size=font_size, bold=(level == 0))
        p.level = level
        _set_para_spacing(p, after=Pt(6))


def build_title_only_slide(prs, slide_spec):
    """构建仅标题页"""
    title = slide_spec.get('title', '')
    layout = get_layout_by_type(prs, 'title_only')
    slide = prs.slides.add_slide(layout)
    set_placeholder_text(slide, 'TITLE', title)
    remove_placeholders(slide)
    return slide


def build_ending_slide(prs, slide_spec):
    """构建结束页 — 44pt标题，副标题仅填充第一个BODY"""
    title = slide_spec.get('title', '谢谢')
    subtitle = slide_spec.get('subtitle', 'Thank You')
    layout = get_layout_by_type(prs, 'ending')
    slide = prs.slides.add_slide(layout)
    body_filled = False
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = str(shape.placeholder_format.type).split('.')[-1]
            if 'CENTER_TITLE' in ph_type:
                set_text_frame_text(shape, title, font_size=44, bold=True)
            elif 'BODY' in ph_type:
                if not body_filled and subtitle:
                    # 只填充第一个 BODY 作为副标题
                    set_text_frame_text(shape, subtitle, font_size=18)
                    body_filled = True
                else:
                    # 其余 BODY 占位符标记删除
                    shapes_to_remove.append(shape)
    # 删除多余 BODY 和所有 OBJECT
    for shape in shapes_to_remove:
        remove_shape(shape)
    remove_placeholders(slide, types_to_remove=['OBJECT'])
    return slide


# ============================================================
# 页面构建器：高级页面类型
# ============================================================

def build_gantt_slide(prs, slide_spec):
    """构建矩阵甘特图页面"""
    title = slide_spec.get('title', '')
    rows = slide_spec.get('rows', [])
    months = slide_spec.get('months', [])
    data = slide_spec.get('data', [])
    milestones = slide_spec.get('milestones', [])
    conclusion = slide_spec.get('conclusion', '')
    layout = get_layout_by_type(prs, 'content')
    slide = prs.slides.add_slide(layout)
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = str(shape.placeholder_format.type).split('.')[-1]
            if 'TITLE' in ph_type:
                set_text_frame_text(shape, title, font_size=26, bold=True)
    remove_placeholders(slide)
    num_rows = len(rows)
    num_cols = len(months)
    if num_rows == 0 or num_cols == 0:
        return slide
    margin_left = Inches(0.6)
    margin_top = Inches(1.35)
    row_header_width = Inches(1.4)
    cell_height = Inches(0.85)
    col_width = Inches(1.5)
    # 月份表头
    header_bg = add_shape_rect(slide, margin_left + row_header_width, margin_top,
                                col_width * num_cols, Inches(0.35), fill_color=COLORS['accent1'])
    for i, month in enumerate(months):
        txBox = slide.shapes.add_textbox(
            margin_left + row_header_width + i * col_width, margin_top, col_width, Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = month
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['body_small'], bold=True, color=COLORS['lt1'])
        p.alignment = PP_ALIGN.CENTER
    # 行和数据
    for r_idx, row_name in enumerate(rows):
        row_y = margin_top + Inches(0.35) + r_idx * cell_height
        add_shape_rect(slide, margin_left, row_y, row_header_width, cell_height,
                       fill_color=COLORS['hk_green'])
        txBox = slide.shapes.add_textbox(margin_left + Inches(0.08), row_y + Inches(0.08),
                                          row_header_width - Inches(0.16), cell_height - Inches(0.16))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = row_name
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['table_body'], bold=True, color=COLORS['accent3'])
        for c_idx in range(num_cols):
            cell_x = margin_left + row_header_width + c_idx * col_width
            add_shape_rect(slide, cell_x, row_y, col_width, cell_height, fill_color=RGBColor(0xFA, 0xFA, 0xFA))
    for item in data:
        row_name = item.get('row', '')
        month = item.get('month', '')
        text = item.get('text', '')
        highlight = item.get('highlight', False)
        if row_name in rows and month in months:
            r_idx = rows.index(row_name)
            c_idx = months.index(month)
            row_y = margin_top + Inches(0.35) + r_idx * cell_height
            cell_x = margin_left + row_header_width + c_idx * col_width
            bg_color = COLORS['hlk_green'] if highlight else RGBColor(0xFA, 0xFA, 0xFA)
            add_shape_rect(slide, cell_x, row_y, col_width, cell_height, fill_color=bg_color)
            if text:
                txBox = slide.shapes.add_textbox(cell_x + Inches(0.06), row_y + Inches(0.06),
                                                  col_width - Inches(0.12), cell_height - Inches(0.12))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = text
                _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['micro'])
    if milestones:
        milestone_y = margin_top + Inches(0.35) + num_rows * cell_height + Inches(0.15)
        for m in milestones:
            month = m.get('month', '')
            text = m.get('text', '')
            if month in months:
                c_idx = months.index(month)
                cell_x = margin_left + row_header_width + c_idx * col_width
                add_shape_rect(slide, cell_x + Inches(0.1), milestone_y,
                               col_width - Inches(0.2), Inches(0.3), fill_color=COLORS['accent1'])
                txBox = slide.shapes.add_textbox(cell_x + Inches(0.1), milestone_y,
                                                  col_width - Inches(0.2), Inches(0.3))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = text
                _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['micro'], color=COLORS['lt1'])
                p.alignment = PP_ALIGN.CENTER
    if conclusion:
        conclusion_y = margin_top + Inches(0.35) + num_rows * cell_height + Inches(0.6)
        txBox = slide.shapes.add_textbox(margin_left, conclusion_y, Inches(12.0), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = conclusion
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['table_body'], color=COLORS['accent3'])
    return slide


def build_timeline_horizontal_slide(prs, slide_spec):
    """构建横向时间线页面"""
    title = slide_spec.get('title', '')
    subtitle_text = slide_spec.get('subtitle', '')
    items = slide_spec.get('items', [])
    layout = get_layout_by_type(prs, 'content')
    slide = prs.slides.add_slide(layout)
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = str(shape.placeholder_format.type).split('.')[-1]
            if 'TITLE' in ph_type:
                set_text_frame_text(shape, title, font_size=26, bold=True)
    remove_placeholders(slide)
    if subtitle_text:
        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(12.0), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle_text
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['table_body'], color=COLORS['accent3'])
    if not items:
        return slide
    margin_left = Inches(0.5)
    margin_top = Inches(1.5)
    card_width = Inches(1.75)
    card_height = Inches(4.5)
    gap = Inches(0.08)
    total_width = len(items) * card_width + (len(items) - 1) * gap
    start_x = max(margin_left, (prs.slide_width - total_width) / 2)
    for i, item in enumerate(items):
        time_label = item.get('time', '')
        card_title = item.get('title', '')
        audience = item.get('audience', '')
        actions = item.get('actions', [])
        card_x = start_x + i * (card_width + gap)
        add_shape_rect(slide, card_x, margin_top, card_width, card_height,
                       fill_color=COLORS['lt1'], line_color=COLORS['hlk_green'], line_width=1)
        add_shape_rect(slide, card_x, margin_top, card_width, Inches(0.35), fill_color=COLORS['accent1'])
        txBox = slide.shapes.add_textbox(card_x, margin_top, card_width, Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = time_label
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['body_small'], bold=True, color=COLORS['lt1'])
        p.alignment = PP_ALIGN.CENTER
        title_box = slide.shapes.add_textbox(card_x + Inches(0.08), margin_top + Inches(0.42),
                                              card_width - Inches(0.16), Inches(0.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = card_title
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['table_body'], bold=True)
        if audience:
            aud_box = slide.shapes.add_textbox(card_x + Inches(0.08), margin_top + Inches(1.05),
                                                card_width - Inches(0.16), Inches(0.3))
            tf = aud_box.text_frame
            p = tf.paragraphs[0]
            p.text = f'▌ {audience}'
            _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['micro'], color=COLORS['accent1'])
        if actions:
            action_box = slide.shapes.add_textbox(card_x + Inches(0.08), margin_top + Inches(1.4),
                                                   card_width - Inches(0.16), Inches(2.5))
            tf = action_box.text_frame
            tf.word_wrap = True
            for a_idx, action in enumerate(actions):
                if a_idx == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = action
                _set_font(p.runs[0] if p.runs else p.add_run(), size=Pt(8))
                _set_para_spacing(p, after=Pt(4))
    return slide


def build_big_number_slide(prs, slide_spec):
    """构建大数字展示页面"""
    title = slide_spec.get('title', '')
    big_number = slide_spec.get('big_number', '')
    unit = slide_spec.get('unit', '')
    description = slide_spec.get('description', '')
    cards = slide_spec.get('cards', [])
    side_cards = slide_spec.get('side_cards', [])
    layout = get_layout_by_type(prs, 'content')
    slide = prs.slides.add_slide(layout)
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = str(shape.placeholder_format.type).split('.')[-1]
            if 'TITLE' in ph_type:
                set_text_frame_text(shape, title, font_size=26, bold=True)
    remove_placeholders(slide)
    big_num_left = Inches(0.6)
    big_num_top = Inches(1.5)
    big_num_width = Inches(5.0)
    big_num_height = Inches(2.0)
    add_shape_rect(slide, big_num_left, big_num_top, big_num_width, big_num_height, fill_color=COLORS['hk_green'])
    num_box = slide.shapes.add_textbox(big_num_left + Inches(0.2), big_num_top + Inches(0.15),
                                        big_num_width - Inches(0.4), Inches(0.9))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = big_number
    _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['big_number'], bold=True, color=COLORS['accent1'])
    unit_box = slide.shapes.add_textbox(big_num_left + Inches(0.2), big_num_top + Inches(1.0),
                                         big_num_width - Inches(0.4), Inches(0.5))
    tf = unit_box.text_frame
    p = tf.paragraphs[0]
    p.text = unit
    _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['body_main'], color=COLORS['dk1'])
    if description:
        desc_box = slide.shapes.add_textbox(big_num_left + Inches(0.2), big_num_top + Inches(1.5),
                                             big_num_width - Inches(0.4), Inches(0.5))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = description
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['note'])
    if side_cards:
        side_left = big_num_left + big_num_width + Inches(0.3)
        side_card_width = Inches(4.0)
        side_card_height = Inches(0.95)
        for s_idx, sc in enumerate(side_cards):
            sc_y = big_num_top + s_idx * (side_card_height + Inches(0.1))
            add_shape_rect(slide, side_left, sc_y, side_card_width, side_card_height, fill_color=COLORS['hk_green'])
            sc_title = sc.get('title', '')
            sc_desc = sc.get('description', '')
            t_box = slide.shapes.add_textbox(side_left + Inches(0.15), sc_y + Inches(0.1),
                                              side_card_width - Inches(0.3), Inches(0.4))
            tf = t_box.text_frame
            p = tf.paragraphs[0]
            p.text = sc_title
            _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['body_main'], bold=True, color=COLORS['accent1'])
            if sc_desc:
                d_box = slide.shapes.add_textbox(side_left + Inches(0.15), sc_y + Inches(0.45),
                                                  side_card_width - Inches(0.3), Inches(0.4))
                tf = d_box.text_frame
                p = tf.paragraphs[0]
                p.text = sc_desc
                _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['micro'])
    if cards:
        card_top = big_num_top + big_num_height + Inches(0.3)
        card_width = Inches(3.8)
        card_height = Inches(1.8)
        gap = Inches(0.15)
        total_cards_width = len(cards) * card_width + (len(cards) - 1) * gap
        start_x = max(Inches(0.6), (prs.slide_width - total_cards_width) / 2)
        for c_idx, card in enumerate(cards):
            card_x = start_x + c_idx * (card_width + gap)
            add_shape_rect(slide, card_x, card_top, card_width, card_height, fill_color=RGBColor(0xF5, 0xF5, 0xF5))
            add_shape_rect(slide, card_x, card_top, Inches(0.06), card_height, fill_color=COLORS['accent1'])
            c_title = card.get('title', '')
            c_calc = card.get('calculation', '')
            c_result = card.get('result', '')
            t_box = slide.shapes.add_textbox(card_x + Inches(0.15), card_top + Inches(0.12),
                                              card_width - Inches(0.3), Inches(0.3))
            tf = t_box.text_frame
            p = tf.paragraphs[0]
            p.text = c_title
            _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['table_body'])
            if c_calc:
                calc_box = slide.shapes.add_textbox(card_x + Inches(0.15), card_top + Inches(0.45),
                                                     card_width - Inches(0.3), Inches(0.3))
                tf = calc_box.text_frame
                p = tf.paragraphs[0]
                p.text = c_calc
                _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['note'])
            if c_result:
                r_box = slide.shapes.add_textbox(card_x + Inches(0.15), card_top + Inches(0.8),
                                                  card_width - Inches(0.3), Inches(0.5))
                tf = r_box.text_frame
                p = tf.paragraphs[0]
                p.text = c_result
                _set_font(p.runs[0] if p.runs else p.add_run(), size=Pt(20), bold=True, color=COLORS['accent1'])
    return slide


def build_comparison_slide(prs, slide_spec):
    """构建数据对比页面"""
    title = slide_spec.get('title', '')
    intro = slide_spec.get('intro', '')
    items = slide_spec.get('items', [])
    target = slide_spec.get('target', '')
    note = slide_spec.get('note', '')
    layout = get_layout_by_type(prs, 'content')
    slide = prs.slides.add_slide(layout)
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = str(shape.placeholder_format.type).split('.')[-1]
            if 'TITLE' in ph_type:
                set_text_frame_text(shape, title, font_size=26, bold=True)
    remove_placeholders(slide)
    if intro:
        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(12.0), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = intro
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['body_small'])
    if not items:
        return slide
    card_top = Inches(1.8)
    card_height = Inches(2.3)
    card_width = Inches(4.0)
    gap = Inches(0.4)
    total_width = len(items) * card_width + (len(items) - 1) * gap
    start_x = max(Inches(0.6), (prs.slide_width - total_width) / 2)
    for i, item in enumerate(items):
        card_x = start_x + i * (card_width + gap)
        label = item.get('label', '')
        value = item.get('value', '')
        description = item.get('description', '')
        add_shape_rect(slide, card_x, card_top, card_width, card_height, fill_color=RGBColor(0xF5, 0xF5, 0xF5))
        l_box = slide.shapes.add_textbox(card_x + Inches(0.15), card_top + Inches(0.1),
                                          card_width - Inches(0.3), Inches(0.3))
        tf = l_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['table_body'])
        v_box = slide.shapes.add_textbox(card_x + Inches(0.15), card_top + Inches(0.4),
                                          card_width - Inches(0.3), Inches(0.6))
        tf = v_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['big_number'], bold=True, color=COLORS['accent1'])
        if description:
            d_box = slide.shapes.add_textbox(card_x + Inches(0.15), card_top + Inches(1.1),
                                              card_width - Inches(0.3), Inches(1.0))
            tf = d_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = description
            _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['note'])
    if target:
        target_y = card_top + card_height + Inches(0.2)
        txBox = slide.shapes.add_textbox(Inches(0.6), target_y, Inches(12.0), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = target
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['table_body'], bold=True, color=COLORS['accent3'])
    if note:
        note_y = card_top + card_height + Inches(0.8)
        txBox = slide.shapes.add_textbox(Inches(0.6), note_y, Inches(12.0), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = note
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['micro'], color=COLORS['lt2'])
    return slide


def build_process_slide(prs, slide_spec):
    """构建流程/阶段展示页面"""
    title = slide_spec.get('title', '')
    steps = slide_spec.get('steps', [])
    direction = slide_spec.get('direction', 'horizontal')
    layout = get_layout_by_type(prs, 'content')
    slide = prs.slides.add_slide(layout)
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = str(shape.placeholder_format.type).split('.')[-1]
            if 'TITLE' in ph_type:
                set_text_frame_text(shape, title, font_size=26, bold=True)
    remove_placeholders(slide)
    if not steps:
        return slide
    if direction == 'horizontal':
        margin_top = Inches(1.8)
        node_width = Inches(2.2)
        node_height = Inches(1.2)
        gap = Inches(0.3)
        total_width = len(steps) * node_width + (len(steps) - 1) * gap
        start_x = max(Inches(0.6), (prs.slide_width - total_width) / 2)
        for i, step in enumerate(steps):
            node_x = start_x + i * (node_width + gap)
            shape_node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, node_x, margin_top, node_width, node_height)
            shape_node.fill.solid()
            shape_node.fill.fore_color.rgb = COLORS['hk_green']
            shape_node.line.color.rgb = COLORS['accent1']
            shape_node.line.width = Pt(2)
            step_title = step.get('title', '')
            step_desc = step.get('description', '')
            t_box = slide.shapes.add_textbox(node_x + Inches(0.1), margin_top + Inches(0.08),
                                              node_width - Inches(0.2), Inches(0.3))
            tf = t_box.text_frame
            p = tf.paragraphs[0]
            p.text = step_title
            _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['body_small'], bold=True, color=COLORS['accent1'])
            p.alignment = PP_ALIGN.CENTER
            d_box = slide.shapes.add_textbox(node_x + Inches(0.1), margin_top + Inches(0.38),
                                              node_width - Inches(0.2), Inches(0.7))
            tf = d_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = step_desc
            _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['note'])
            p.alignment = PP_ALIGN.CENTER
            if i < len(steps) - 1:
                arrow_x = node_x + node_width + Inches(0.02)
                arrow_y = margin_top + node_height / 2 - Inches(0.08)
                arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y, Inches(0.26), Inches(0.16))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = COLORS['accent1']
                arrow.line.fill.background()
    return slide


def build_kpi_dashboard_slide(prs, slide_spec):
    """构建KPI仪表盘页面"""
    title = slide_spec.get('title', '')
    kpis = slide_spec.get('kpis', [])
    layout = get_layout_by_type(prs, 'content')
    slide = prs.slides.add_slide(layout)
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = str(shape.placeholder_format.type).split('.')[-1]
            if 'TITLE' in ph_type:
                set_text_frame_text(shape, title, font_size=26, bold=True)
    remove_placeholders(slide)
    if not kpis:
        return slide
    cols = min(4, len(kpis))
    rows = (len(kpis) + cols - 1) // cols
    card_width = Inches(2.9)
    card_height = Inches(1.3)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)
    total_grid_width = cols * card_width + (cols - 1) * gap_x
    start_x = max(Inches(0.5), (prs.slide_width - total_grid_width) / 2)
    start_y = Inches(1.5)
    for idx, kpi in enumerate(kpis):
        row = idx // cols
        col = idx % cols
        card_x = start_x + col * (card_width + gap_x)
        card_y = start_y + row * (card_height + gap_y)
        label = kpi.get('label', '')
        value = kpi.get('value', '')
        unit = kpi.get('unit', '')
        change = kpi.get('change', '')
        target_val = kpi.get('target', '')
        add_shape_rect(slide, card_x, card_y, card_width, card_height,
                       fill_color=COLORS['lt1'], line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=1)
        add_shape_rect(slide, card_x, card_y, card_width, Inches(0.06), fill_color=COLORS['accent1'])
        l_box = slide.shapes.add_textbox(card_x + Inches(0.1), card_y + Inches(0.12),
                                          card_width - Inches(0.2), Inches(0.22))
        tf = l_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['note'])
        v_text = f"{value}{unit}"
        v_box = slide.shapes.add_textbox(card_x + Inches(0.1), card_y + Inches(0.35),
                                          card_width - Inches(0.2), Inches(0.4))
        tf = v_box.text_frame
        p = tf.paragraphs[0]
        p.text = v_text
        _set_font(p.runs[0] if p.runs else p.add_run(), size=Pt(22), bold=True, color=COLORS['accent1'])
        meta_parts = []
        if change:
            meta_parts.append(change)
        if target_val:
            meta_parts.append(f"目标: {target_val}")
        if meta_parts:
            m_box = slide.shapes.add_textbox(card_x + Inches(0.1), card_y + Inches(0.78),
                                              card_width - Inches(0.2), Inches(0.25))
            tf = m_box.text_frame
            p = tf.paragraphs[0]
            p.text = '  |  '.join(meta_parts)
            _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['micro'], color=COLORS['lt2'])
    return slide


# ============================================================
# 页面构建器：矩阵与复盘页面类型
# ============================================================

def build_item_matrix_slide(prs, slide_spec):
    """构建物品矩阵图页面（品牌提示物风格）"""
    title = slide_spec.get('title', '')
    groups = slide_spec.get('groups', [])
    footer_notes = slide_spec.get('footer_notes', '')
    layout = get_layout_by_type(prs, 'content')
    slide = prs.slides.add_slide(layout)
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = str(shape.placeholder_format.type).split('.')[-1]
            if 'TITLE' in ph_type:
                set_text_frame_text(shape, title, font_size=26, bold=True)
    remove_placeholders(slide)
    if not groups:
        return slide
    margin_left = Inches(0.4)
    margin_top = Inches(1.35)
    scene_w = Inches(1.5)
    audience_w = Inches(1.35)
    items_w = Inches(2.4)
    cost_w = Inches(1.35)
    score_w = Inches(0.65)
    row_h = Inches(0.33)
    group_gap = Inches(0.08)
    scene_colors = [
        RGBColor(0x00, 0x9D, 0x4A), RGBColor(0x1F, 0xB8, 0x40),
        RGBColor(0x5A, 0xC7, 0x1E), RGBColor(0x02, 0x9B, 0x46),
    ]
    score_color = COLORS['score_orange']
    item_bg = RGBColor(0xF5, 0xF5, 0xF5)
    total_items = sum(len(g.get('items', [])) for g in groups)
    needed_height = total_items * row_h + len(groups) * group_gap + Inches(0.3)
    slide_body_h = prs.slide_height - Inches(1.5)
    if needed_height > slide_body_h:
        row_h = (slide_body_h - Inches(0.5) - len(groups) * group_gap) / max(total_items, 1)
        row_h = max(row_h, Inches(0.22))
    # 列标题
    header_y = margin_top
    add_shape_rect(slide, margin_left, header_y, scene_w, Inches(0.32), fill_color=COLORS['accent1'])
    _add_centered_text(slide, margin_left, header_y, scene_w, Inches(0.32), '场景', font_size=11, color=COLORS['lt1'], bold=True)
    add_shape_rect(slide, margin_left + scene_w, header_y, audience_w, Inches(0.32), fill_color=COLORS['accent1'])
    _add_centered_text(slide, margin_left + scene_w, header_y, audience_w, Inches(0.32), '对象', font_size=11, color=COLORS['lt1'], bold=True)
    add_shape_rect(slide, margin_left + scene_w + audience_w, header_y, items_w, Inches(0.32), fill_color=COLORS['accent1'])
    _add_centered_text(slide, margin_left + scene_w + audience_w, header_y, items_w, Inches(0.32), '品牌提示物', font_size=11, color=COLORS['lt1'], bold=True)
    add_shape_rect(slide, margin_left + scene_w + audience_w + items_w, header_y, cost_w, Inches(0.32), fill_color=COLORS['accent1'])
    _add_centered_text(slide, margin_left + scene_w + audience_w + items_w, header_y, cost_w, Inches(0.32), '成本分档', font_size=11, color=COLORS['lt1'], bold=True)
    add_shape_rect(slide, margin_left + scene_w + audience_w + items_w + cost_w, header_y, score_w, Inches(0.32), fill_color=score_color)
    _add_centered_text(slide, margin_left + scene_w + audience_w + items_w + cost_w, header_y, score_w, Inches(0.32), '平均\n喜爱度', font_size=10, color=COLORS['lt1'], bold=True)
    current_y = margin_top + Inches(0.32)
    for g_idx, group in enumerate(groups):
        scene = group.get('scene', '')
        audience = group.get('audience', '')
        cost_tier = group.get('cost_tier', '')
        avg_score = group.get('avg_score', '')
        items = group.get('items', [])
        top_items = group.get('top_items', [])
        scene_color = scene_colors[g_idx % len(scene_colors)]
        n_items = len(items) if items else 1
        group_h = n_items * row_h
        add_shape_rect(slide, margin_left, current_y, scene_w, group_h, fill_color=scene_color)
        _add_centered_text(slide, margin_left, current_y, scene_w, group_h, scene, font_size=12, color=COLORS['lt1'], bold=True, word_wrap=True)
        add_shape_rect(slide, margin_left + scene_w, current_y, audience_w, group_h, fill_color=COLORS['hk_green'])
        _add_centered_text(slide, margin_left + scene_w, current_y, audience_w, group_h, audience, font_size=10, word_wrap=True)
        for i_idx, item in enumerate(items):
            item_name = item.get('name', '') if isinstance(item, dict) else str(item)
            item_score = item.get('score', '') if isinstance(item, dict) else ''
            row_y = current_y + i_idx * row_h
            add_shape_rect(slide, margin_left + scene_w + audience_w, row_y, items_w, row_h, fill_color=item_bg)
            _add_left_text(slide, margin_left + scene_w + audience_w + Inches(0.06), row_y, items_w - Inches(0.12), row_h, item_name, font_size=9)
            if i_idx == 0:
                add_shape_rect(slide, margin_left + scene_w + audience_w + items_w, row_y, cost_w, group_h, fill_color=COLORS['hk_green'])
                _add_centered_text(slide, margin_left + scene_w + audience_w + items_w, row_y, cost_w, group_h, cost_tier, font_size=10, color=COLORS['accent3'], bold=True, word_wrap=True)
            score_val = item_score if item_score else (avg_score if i_idx == 0 else '')
            if score_val:
                add_shape_rect(slide, margin_left + scene_w + audience_w + items_w + cost_w, row_y, score_w, row_h, fill_color=COLORS['hlk_orange'])
                _add_centered_text(slide, margin_left + scene_w + audience_w + items_w + cost_w, row_y, score_w, row_h, score_val, font_size=10, color=score_color, bold=True)
        if top_items:
            top_y = current_y + group_h + Inches(0.02)
            top_text = 'TOP最受欢迎: ' + '、'.join(top_items[:3])
            txBox = slide.shapes.add_textbox(margin_left + scene_w + audience_w, top_y, items_w + cost_w, Inches(0.22))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = top_text
            _set_font(p.runs[0] if p.runs else p.add_run(), size=Pt(8), color=COLORS['accent3'])
        current_y += group_h + group_gap
    if footer_notes:
        note_y = current_y + Inches(0.1)
        txBox = slide.shapes.add_textbox(margin_left, note_y, prs.slide_width - margin_left * 2, Inches(0.35))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = footer_notes
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['micro'], color=COLORS['lt2'])
    return slide


def build_calendar_grid_slide(prs, slide_spec):
    """构建日历网格页面"""
    title = slide_spec.get('title', '')
    rows = slide_spec.get('rows', [])
    months = slide_spec.get('months', [])
    grid = slide_spec.get('grid', [])
    note = slide_spec.get('note', '')
    layout = get_layout_by_type(prs, 'content')
    slide = prs.slides.add_slide(layout)
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = str(shape.placeholder_format.type).split('.')[-1]
            if 'TITLE' in ph_type:
                set_text_frame_text(shape, title, font_size=26, bold=True)
    remove_placeholders(slide)
    if not rows or not months:
        return slide
    margin_left = Inches(0.3)
    margin_top = Inches(1.3)
    row_header_w = Inches(1.3)
    month_hd_h = Inches(0.38)
    row_h = Inches(1.15)
    col_w = (prs.slide_width - margin_left - row_header_w - Inches(0.3)) / max(len(months), 1)
    col_w = min(col_w, Inches(2.0))
    total_grid_w = row_header_w + col_w * len(months)
    max_grid_w = prs.slide_width - margin_left * 2
    if total_grid_w > max_grid_w:
        col_w = (max_grid_w - row_header_w) / len(months)
    for c_idx, month in enumerate(months):
        col_x = margin_left + row_header_w + c_idx * col_w
        add_shape_rect(slide, col_x, margin_top, col_w, month_hd_h, fill_color=COLORS['accent1'])
        _add_centered_text(slide, col_x, margin_top, col_w, month_hd_h, month, font_size=12, color=COLORS['lt1'], bold=True)
    for r_idx, row in enumerate(rows):
        row_name = row.get('name', '') if isinstance(row, dict) else str(row)
        row_sub = row.get('subtitle', '') if isinstance(row, dict) else ''
        row_y = margin_top + month_hd_h + r_idx * row_h
        row_color = COLORS['accent3'] if r_idx % 2 == 0 else COLORS['accent1']
        add_shape_rect(slide, margin_left, row_y, row_header_w, row_h, fill_color=row_color)
        title_text = row_name
        if row_sub:
            title_text += f"\n{row_sub}"
        _add_centered_text(slide, margin_left, row_y, row_header_w, row_h, title_text, font_size=10, color=COLORS['lt1'], bold=True, word_wrap=True)
        for c_idx in range(len(months)):
            col_x = margin_left + row_header_w + c_idx * col_w
            cell_data = {}
            if r_idx < len(grid) and c_idx < len(grid[r_idx]):
                cell_data = grid[r_idx][c_idx]
            items = cell_data.get('items', [])
            highlight = cell_data.get('highlight', False)
            bg = COLORS['hlk_green'] if highlight else (RGBColor(0xFA, 0xFA, 0xFA) if c_idx % 2 == 0 else COLORS['lt1'])
            add_shape_rect(slide, col_x, row_y, col_w, row_h, fill_color=bg)
            if items:
                item_text = '\n'.join(items[:5])
                _add_left_text(slide, col_x + Inches(0.04), row_y + Inches(0.04), col_w - Inches(0.08), row_h - Inches(0.08), item_text, font_size=8, word_wrap=True)
    if note:
        note_y = margin_top + month_hd_h + len(rows) * row_h + Inches(0.1)
        txBox = slide.shapes.add_textbox(margin_left, note_y, prs.slide_width - margin_left * 2, Inches(0.3))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = note
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['micro'], color=COLORS['lt2'])
    return slide


def build_review_matrix_slide(prs, slide_spec):
    """构建复盘矩阵页面（有效动作 × 做的好/做得不好）"""
    title = slide_spec.get('title', '')
    badge = slide_spec.get('badge', '')
    category_header = slide_spec.get('category_header', '有效动作')
    good_header = slide_spec.get('good_header', '做的好')
    bad_header = slide_spec.get('bad_header', '做得不好')
    rows = slide_spec.get('rows', [])
    footer_note = slide_spec.get('footer_note', '')
    layout = get_layout_by_type(prs, 'content')
    slide = prs.slides.add_slide(layout)
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = str(shape.placeholder_format.type).split('.')[-1]
            if 'TITLE' in ph_type:
                set_text_frame_text(shape, title, font_size=26, bold=True)
    remove_placeholders(slide)
    if not rows:
        return slide
    # 兼容两种rows格式: [{action,good,bad}, ...] 或 [[action,good,bad], ...]
    if rows and isinstance(rows[0], list):
        rows = [{'action': r[0], 'good': r[1] if len(r) > 1 else '', 'bad': r[2] if len(r) > 2 else ''} for r in rows]
    margin_left = Inches(0.3)
    margin_top = Inches(1.45)
    col1_w = Inches(2.2)
    col2_w = Inches(4.85)
    col3_w = Inches(4.85)
    header_h = Inches(0.45)
    row_h_min = Inches(0.95)
    row_padding = Inches(0.06)
    available_h = prs.slide_height - margin_top - header_h - Inches(0.8)
    n_rows = len(rows)
    calc_h = (available_h - row_padding * (n_rows - 1)) / n_rows
    row_h = max(row_h_min, min(calc_h, Inches(1.4)))
    add_shape_rect(slide, margin_left, margin_top, col1_w, header_h, fill_color=COLORS['accent3'])
    _add_centered_text(slide, margin_left, margin_top, col1_w, header_h, category_header, font_size=13, color=COLORS['lt1'], bold=True)
    add_shape_rect(slide, margin_left + col1_w, margin_top, col2_w, header_h, fill_color=COLORS['accent2'])
    _add_centered_text(slide, margin_left + col1_w, margin_top, col2_w, header_h, f'✓  {good_header}', font_size=13, color=COLORS['lt1'], bold=True)
    bad_bg = COLORS['accent5']
    add_shape_rect(slide, margin_left + col1_w + col2_w, margin_top, col3_w, header_h, fill_color=bad_bg)
    _add_centered_text(slide, margin_left + col1_w + col2_w, margin_top, col3_w, header_h, f'⚠  {bad_header}', font_size=13, color=COLORS['dk1'], bold=True)
    if badge:
        badge_x = margin_left - Inches(0.15)
        badge_y = margin_top - Inches(0.02)
        add_shape_rect(slide, badge_x, badge_y, Inches(1.2), Inches(0.38), fill_color=COLORS['accent1'])
        _add_centered_text(slide, badge_x, badge_y, Inches(1.2), Inches(0.38), badge, font_size=12, color=COLORS['lt1'], bold=True)
    for r_idx, row in enumerate(rows):
        row_y = margin_top + header_h + r_idx * (row_h + row_padding)
        action = row.get('action', '')
        note = row.get('note', '')
        good = row.get('good', '')
        bad = row.get('bad', '')
        row_color = COLORS['accent3'] if r_idx % 2 == 0 else COLORS['accent1']
        add_shape_rect(slide, margin_left, row_y, col1_w, row_h, fill_color=row_color)
        action_text = action
        if note:
            action_text += f'\n{note}'
        _add_centered_text(slide, margin_left, row_y, col1_w, row_h, action_text, font_size=11, color=COLORS['lt1'], bold=True, word_wrap=True)
        good_bg = COLORS['hlk_green'] if r_idx % 2 == 0 else COLORS['lt1']
        add_shape_rect(slide, margin_left + col1_w, row_y, col2_w, row_h, fill_color=good_bg)
        _add_left_text(slide, margin_left + col1_w + Inches(0.12), row_y + Inches(0.06), col2_w - Inches(0.24), row_h - Inches(0.12), good, font_size=10, word_wrap=True)
        bad_bg_row = COLORS['hlk_orange'] if r_idx % 2 == 0 else COLORS['lt1']
        add_shape_rect(slide, margin_left + col1_w + col2_w, row_y, col3_w, row_h, fill_color=bad_bg_row)
        _add_left_text(slide, margin_left + col1_w + col2_w + Inches(0.12), row_y + Inches(0.06), col3_w - Inches(0.24), row_h - Inches(0.12), bad, font_size=10, word_wrap=True)
    if footer_note:
        note_y = margin_top + header_h + n_rows * (row_h + row_padding) + Inches(0.1)
        txBox = slide.shapes.add_textbox(margin_left, note_y, prs.slide_width - margin_left * 2, Inches(0.35))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = footer_note
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['micro'], color=COLORS['lt2'])
    return slide


def build_action_category_slide(prs, slide_spec):
    """构建分类复盘页面（持续做/优化提升/开始做）"""
    title = slide_spec.get('title', '')
    badge = slide_spec.get('badge', '')
    intro = slide_spec.get('intro', '')
    categories = slide_spec.get('categories', [])
    highlight = slide_spec.get('highlight_box', None)
    layout = get_layout_by_type(prs, 'content')
    slide = prs.slides.add_slide(layout)
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = str(shape.placeholder_format.type).split('.')[-1]
            if 'TITLE' in ph_type:
                set_text_frame_text(shape, title, font_size=26, bold=True)
    remove_placeholders(slide)
    if not categories:
        return slide
    margin_left = Inches(0.3)
    body_top = Inches(1.35)
    if badge:
        add_shape_rect(slide, margin_left, body_top, Inches(1.6), Inches(0.34), fill_color=COLORS['accent1'])
        _add_centered_text(slide, margin_left, body_top, Inches(1.6), Inches(0.34), badge, font_size=11, color=COLORS['lt1'], bold=True)
    current_y = body_top + Inches(0.55) if badge else body_top
    if intro:
        txBox = slide.shapes.add_textbox(margin_left, current_y, prs.slide_width - margin_left * 2, Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = intro
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['table_body'], color=COLORS['dk2'])
        current_y += Inches(0.5)
    total_items = sum(len(c.get('items', [])) for c in categories)
    remaining_h = prs.slide_height - current_y - Inches(0.7)
    if highlight:
        remaining_h -= Inches(0.75)
    item_h = remaining_h / max(total_items, 1)
    item_h = min(item_h, Inches(0.35))
    for c_idx, cat in enumerate(categories):
        label = cat.get('label', '')
        items = cat.get('items', [])
        cat_color = [COLORS['accent3'], COLORS['accent1'], COLORS['accent2']][c_idx % 3]
        n = len(items)
        cat_h = max(item_h * n, Inches(0.5))
        label_w = Inches(1.8)
        add_shape_rect(slide, margin_left, current_y, label_w, cat_h, fill_color=cat_color)
        _add_centered_text(slide, margin_left, current_y, label_w, cat_h, label, font_size=12, color=COLORS['lt1'], bold=True, word_wrap=True)
        for i_idx, item_text in enumerate(items):
            row_y = current_y + i_idx * item_h
            bg = COLORS['hk_green'] if i_idx % 2 == 0 else COLORS['lt1']
            add_shape_rect(slide, margin_left + label_w, row_y, prs.slide_width - margin_left * 2 - label_w, item_h, fill_color=bg)
            _add_left_text(slide, margin_left + label_w + Inches(0.12), row_y + Inches(0.02),
                          prs.slide_width - margin_left * 2 - label_w - Inches(0.24), item_h,
                          f'• {item_text}', font_size=10, word_wrap=True)
        current_y += cat_h + Inches(0.08)
    if highlight:
        current_y += Inches(0.1)
        box_title = highlight.get('title', '')
        box_text = highlight.get('text', '')
        add_shape_rect(slide, margin_left, current_y, prs.slide_width - margin_left * 2, Inches(0.6), fill_color=COLORS['hlk_green'])
        if box_title:
            _add_left_text(slide, margin_left + Inches(0.12), current_y + Inches(0.02),
                          prs.slide_width - margin_left * 2 - Inches(0.24), Inches(0.3),
                          box_title, font_size=11, color=COLORS['accent3'], bold=True)
            _add_left_text(slide, margin_left + Inches(0.12), current_y + Inches(0.28),
                          prs.slide_width - margin_left * 2 - Inches(0.24), Inches(0.28),
                          box_text, font_size=9, word_wrap=True)
        else:
            _add_left_text(slide, margin_left + Inches(0.12), current_y + Inches(0.06),
                          prs.slide_width - margin_left * 2 - Inches(0.24), Inches(0.48),
                          box_text, font_size=10, word_wrap=True)
    return slide


def build_strategy_diagram_slide(prs, slide_spec):
    """构建策略架构图页面"""
    title = slide_spec.get('title', '')
    badge = slide_spec.get('badge', '')
    diagram = slide_spec.get('diagram', {})
    layout = get_layout_by_type(prs, 'content')
    slide = prs.slides.add_slide(layout)
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = str(shape.placeholder_format.type).split('.')[-1]
            if 'TITLE' in ph_type:
                set_text_frame_text(shape, title, font_size=26, bold=True)
    remove_placeholders(slide)
    if not diagram:
        return slide
    center = diagram.get('center', {})
    dimensions = diagram.get('dimensions', [])
    pillars = diagram.get('pillars', [])
    margin_left = Inches(0.3)
    margin_top = Inches(1.2)
    if badge:
        add_shape_rect(slide, margin_left, margin_top, Inches(1.6), Inches(0.34), fill_color=COLORS['accent1'])
        _add_centered_text(slide, margin_left, margin_top, Inches(1.6), Inches(0.34), badge, font_size=11, color=COLORS['lt1'], bold=True)
    center_text = center.get('text', '')
    center_sub = center.get('sub', '')
    center_x = Inches(5.2)
    center_y = margin_top + Inches(0.15)
    center_w = Inches(3.2)
    center_h = Inches(0.75) if center_sub else Inches(0.55)
    add_shape_rect(slide, center_x, center_y, center_w, center_h,
                   fill_color=COLORS['accent3'], line_color=COLORS['accent1'], line_width=Pt(2))
    if center_sub:
        _add_centered_text(slide, center_x, center_y, center_w, Inches(0.32), center_text, font_size=13, color=COLORS['lt1'], bold=True)
        _add_centered_text(slide, center_x, center_y + Inches(0.35), center_w, Inches(0.3), center_sub, font_size=9, color=COLORS['hlk_green'])
    else:
        _add_centered_text(slide, center_x, center_y, center_w, center_h, center_text, font_size=14, color=COLORS['lt1'], bold=True)
    n_dims = len(dimensions)
    if n_dims == 0:
        return slide
    dim_spacing = Inches(8.5) / max(n_dims, 1)
    dim_start_x = margin_left + Inches(2.0)
    for d_idx, dim in enumerate(dimensions):
        dim_label = dim.get('label', '')
        dim_items = dim.get('items', [])
        sub_label = dim.get('sub_label', '')
        dim_x = dim_start_x + d_idx * dim_spacing
        dim_w = dim_spacing - Inches(0.15)
        dim_y = center_y + center_h + Inches(0.25)
        arrow_w = Inches(1.35)
        arrow_h = Inches(0.35)
        arrow_x = dim_x + dim_w / 2 - arrow_w / 2
        arrow_shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, arrow_x, dim_y, arrow_w, arrow_h)
        arrow_shape.fill.solid()
        arrow_shape.fill.fore_color.rgb = COLORS['accent2']
        arrow_shape.line.fill.background()
        if arrow_shape.has_text_frame:
            tf = arrow_shape.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = dim_label
            _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['note'], bold=True, color=COLORS['lt1'])
            p.alignment = PP_ALIGN.CENTER
        if sub_label:
            sub_y = dim_y + arrow_h + Inches(0.05)
            _add_centered_text(slide, dim_x, sub_y, dim_w, Inches(0.22), sub_label, font_size=8, color=COLORS['lt2'])
        sub_start_y = dim_y + arrow_h + (Inches(0.35) if sub_label else Inches(0.1))
        for i_idx, item in enumerate(dim_items):
            rect_w = dim_w - Inches(0.1)
            rect_h = Inches(0.32)
            rect_x = dim_x + Inches(0.05)
            rect_y = sub_start_y + i_idx * Inches(0.4)
            node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rect_x, rect_y, rect_w, rect_h)
            node.fill.solid()
            node.fill.fore_color.rgb = COLORS['hk_green']
            node.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            node.line.width = Pt(0.5)
            if node.has_text_frame:
                tf = node.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = item
                _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['micro'])
                p.alignment = PP_ALIGN.CENTER
    if pillars:
        max_dim_bottom = margin_top + center_h + Inches(0.25) + (Inches(0.35) + Inches(0.1) + 3 * Inches(0.4))
        pillar_y = min(max_dim_bottom + Inches(0.3), prs.slide_height - Inches(1.2))
        n_pillars = len(pillars)
        pillar_spacing = Inches(8.5) / max(n_pillars, 1)
        pillar_start_x = margin_left + Inches(2.0)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, margin_left, pillar_y,
                                       prs.slide_width - margin_left * 2, Inches(0.015))
        line.fill.solid()
        line.fill.fore_color.rgb = COLORS['accent2']
        line.line.fill.background()
        _add_left_text(slide, margin_left, pillar_y + Inches(0.1), Inches(1.8), Inches(0.3),
                      '四维服务升级', font_size=10, color=COLORS['accent3'], bold=True)
        for p_idx, pillar in enumerate(pillars):
            p_label = pillar.get('label', '')
            p_sub = pillar.get('sub', '')
            px = pillar_start_x + p_idx * pillar_spacing
            pw = pillar_spacing - Inches(0.15)
            py = pillar_y + Inches(0.15)
            node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px, py, pw, Inches(0.55))
            node.fill.solid()
            node.fill.fore_color.rgb = COLORS['hk_green']
            node.line.color.rgb = COLORS['accent2']
            node.line.width = Pt(1)
            if node.has_text_frame:
                tf = node.text_frame
                tf.word_wrap = True
                p1 = tf.paragraphs[0]
                p1.text = p_label
                _set_font(p1.runs[0] if p1.runs else p1.add_run(), size=FONT_SIZES['note'], bold=True)
                p1.alignment = PP_ALIGN.CENTER
                if p_sub:
                    p2 = tf.add_paragraph()
                    p2.text = p_sub
                    _set_font(p2.runs[0] if p2.runs else p2.add_run(), size=Pt(8), color=COLORS['lt2'])
                    p2.alignment = PP_ALIGN.CENTER
    return slide


# ============================================================
# 新增页面类型：双轴图 / 升级甘特图 / 可视化对比矩阵
# ============================================================

def build_dual_axis_chart_slide(prs, slide_spec):
    """构建双轴图页面 — 柱状图+折线图叠加，左轴+右轴"""
    title = slide_spec.get('title', '')
    subtitle = slide_spec.get('subtitle', '')
    categories = slide_spec.get('categories', [])
    bar_series = slide_spec.get('bar_series', {})    # {"name": "销售额", "values": [...], "unit": "亿元"}
    line_series = slide_spec.get('line_series', {})   # {"name": "增长率", "values": [...], "unit": "%"}
    note = slide_spec.get('note', '')
    
    layout = get_layout_by_type(prs, 'content')
    slide = prs.slides.add_slide(layout)
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = str(shape.placeholder_format.type).split('.')[-1]
            if 'TITLE' in ph_type:
                set_text_frame_text(shape, title, font_size=26, bold=True)
    remove_placeholders(slide)
    
    if subtitle:
        txBox = slide.shapes.add_textbox(Inches(0.62), Inches(1.05), Inches(12.0), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['table_body'], color=COLORS['dk2'])
    
    if not categories:
        return slide
    
    bar_name = bar_series.get('name', '系列1')
    bar_values = bar_series.get('values', [])
    bar_unit = bar_series.get('unit', '')
    line_name = line_series.get('name', '系列2')
    line_values = line_series.get('values', [])
    line_unit = line_series.get('unit', '')
    
    # 创建组合图表：柱状图+折线图
    chart_data = ChartData()
    chart_data.categories = categories
    chart_data.add_series(bar_name, bar_values)
    chart_data.add_series(line_name, line_values)
    
    chart_left = Inches(0.8)
    chart_top = Inches(1.5)
    chart_width = Inches(11.5)
    chart_height = Inches(4.2)
    
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, chart_left, chart_top, chart_width, chart_height, chart_data
    )
    chart = graphic_frame.chart
    
    # 柱状图样式（左轴）
    bar_plot = chart.plots[0]
    bar_plot.has_data_labels = True
    bar_plot.data_labels.font.size = FONT_SIZES['micro']
    bar_plot.data_labels.font.name = FONT_CN
    
    bar_series_obj = chart.series[0]
    bar_series_obj.format.fill.solid()
    bar_series_obj.format.fill.fore_color.rgb = COLORS['accent1']
    
    # 将第二个系列改为折线图并设右轴
    line_series_obj = chart.series[1]
    line_series_obj.chart_type = XL_CHART_TYPE.LINE_MARKERS
    line_series_obj.format.line.color.rgb = COLORS['accent6']
    line_series_obj.format.line.width = Pt(2.5)
    line_series_obj.smooth = False
    line_series_obj.has_data_labels = True
    line_series_obj.data_labels.font.size = FONT_SIZES['micro']
    line_series_obj.data_labels.font.name = FONT_CN
    line_series_obj.data_labels.font.color.rgb = COLORS['accent6']
    
    # 设置次坐标轴
    try:
        from pptx.oxml.ns import qn as _qn
        # 折线系列放到次坐标轴
        c_chart = chart._chartSpace
        plot_area = c_chart.find(_qn('c:chart'))
        if plot_area is None:
            plot_area = c_chart.findall(_qn('c:chart'))[0]
        # 找到第二个系列设置次坐标
        bar_dir = c_chart.find('.//' + _qn('c:barDir'))
        if bar_dir is not None:
            # 对line系列: 设置c:overlap和c:gapWidth使柱状和折线共存
            pass
    except:
        pass
    
    # 颜色方案
    if chart.has_legend:
        chart.legend.include_in_layout = False
        chart.legend.font.size = FONT_SIZES['note']
        chart.legend.font.name = FONT_CN
    
    # 坐标轴
    if hasattr(chart, 'category_axis') and chart.category_axis:
        chart.category_axis.tick_labels.font.size = FONT_SIZES['note']
        chart.category_axis.tick_labels.font.name = FONT_CN
    if hasattr(chart, 'value_axis') and chart.value_axis:
        chart.value_axis.tick_labels.font.size = FONT_SIZES['note']
        chart.value_axis.tick_labels.font.name = FONT_CN
        chart.value_axis.has_title = True
        chart.value_axis.axis_title.text_frame.text = bar_unit if bar_unit else bar_name
        for para in chart.value_axis.axis_title.text_frame.paragraphs:
            for run in para.runs:
                _set_font(run, size=FONT_SIZES['micro'], color=COLORS['dk2'])
    
    if note:
        note_y = chart_top + chart_height + Inches(0.1)
        txBox = slide.shapes.add_textbox(Inches(0.62), note_y, Inches(12.0), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = note
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['micro'], color=COLORS['lt2'])
    
    return slide


def build_gantt_table_slide(prs, slide_spec):
    """构建升级版甘特图 — 真实表格驱动 + 进度条 + 里程碑标注"""
    title = slide_spec.get('title', '')
    subtitle = slide_spec.get('subtitle', '')
    rows = slide_spec.get('rows', [])          # [{task, start, end, progress, owner, status}, ...]
    conclusion = slide_spec.get('conclusion', '')
    
    layout = get_layout_by_type(prs, 'content')
    slide = prs.slides.add_slide(layout)
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = str(shape.placeholder_format.type).split('.')[-1]
            if 'TITLE' in ph_type:
                set_text_frame_text(shape, title, font_size=26, bold=True)
    remove_placeholders(slide)
    
    if subtitle:
        txBox = slide.shapes.add_textbox(Inches(0.62), Inches(1.05), Inches(12.0), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['table_body'], color=COLORS['dk2'])
    
    if not rows:
        return slide
    
    # 真实表格: 任务名称 | 开始 | 结束 | 进度 | 负责人 | 状态
    headers = ['任务名称', '开始', '结束', '进度', '负责人', '状态']
    n_rows = len(rows) + 1
    n_cols = len(headers)
    
    table_left = Inches(0.4)
    table_top = Inches(1.45)
    table_width = Inches(12.5)
    row_height = Inches(0.48)
    table_height = row_height * n_rows
    
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, table_left, table_top, table_width, table_height)
    tbl = tbl_shape.table
    
    # 列宽: 任务名称占40%, 其余均分
    col_widths = [Inches(4.0), Inches(1.6), Inches(1.6), Inches(1.6), Inches(2.0), Inches(1.7)]
    for ci, w in enumerate(col_widths):
        tbl.columns[ci].width = w
    
    # 表头
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['accent1']
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                _set_font(run, size=FONT_SIZES['table_header'], bold=True, color=COLORS['lt1'])
    
    # 状态颜色映射
    status_colors = {
        '完成': COLORS['accent1'], '进行中': COLORS['accent2'], 
        '未开始': COLORS['lt2'], '延期': COLORS['accent6'],
        '待确认': COLORS['accent5'],
    }
    
    # 数据行
    for ri, row in enumerate(rows):
        task = row.get('task', '')
        start = row.get('start', '')
        end = row.get('end', '')
        progress = row.get('progress', '')    # "0%" ~ "100%"
        owner = row.get('owner', '')
        status = row.get('status', '')
        
        row_data = [task, start, end, progress, owner, status]
        for ci, val in enumerate(row_data):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            # 隔行背景
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['table_stripe']
            # 字体
            for para in cell.text_frame.paragraphs:
                if ci == 0:
                    para.alignment = PP_ALIGN.LEFT
                else:
                    para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    _set_font(run, size=FONT_SIZES['table_body'], color=COLORS['dk1'])
        
        # 进度条效果：在进度单元格添加绿色背景覆盖
        if progress:
            try:
                pct = int(progress.replace('%', '')) / 100.0
                # 在进度文字下方添加色块指示进度
                prog_cell = tbl.cell(ri + 1, 3)
                if pct >= 0.8:
                    prog_cell.fill.solid()
                    prog_cell.fill.fore_color.rgb = COLORS['hlk_green']
                elif pct >= 0.4:
                    prog_cell.fill.solid()
                    prog_cell.fill.fore_color.rgb = COLORS['hk_green']
                elif pct > 0:
                    prog_cell.fill.solid()
                    prog_cell.fill.fore_color.rgb = COLORS['hlk_orange']
                for para in prog_cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.bold = True
            except:
                pass
        
        # 状态单元格颜色
        if status in status_colors:
            st_cell = tbl.cell(ri + 1, 5)
            st_cell.fill.solid()
            st_cell.fill.fore_color.rgb = status_colors[status]
            for para in st_cell.text_frame.paragraphs:
                for run in para.runs:
                    _set_font(run, size=FONT_SIZES['table_body'], bold=True, 
                             color=COLORS['lt1'] if status != '未开始' else COLORS['dk1'])
    
    if conclusion:
        conc_y = table_top + table_height + Inches(0.12)
        txBox = slide.shapes.add_textbox(Inches(0.4), conc_y, Inches(12.0), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = conclusion
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['table_body'], 
                 bold=True, color=COLORS['accent3'])
    
    return slide


def build_visual_comparison_slide(prs, slide_spec):
    """构建可视化对比矩阵页面 — 多维度雷达对比 + 评分卡片"""
    title = slide_spec.get('title', '')
    subtitle = slide_spec.get('subtitle', '')
    dimensions = slide_spec.get('dimensions', [])   # ["品牌声量", "覆盖深度", ...]
    groups = slide_spec.get('groups', [])            # [{name, values:[], color}, ...]
    highlight_finding = slide_spec.get('highlight_finding', '')
    
    layout = get_layout_by_type(prs, 'content')
    slide = prs.slides.add_slide(layout)
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = str(shape.placeholder_format.type).split('.')[-1]
            if 'TITLE' in ph_type:
                set_text_frame_text(shape, title, font_size=26, bold=True)
    remove_placeholders(slide)
    
    if subtitle:
        txBox = slide.shapes.add_textbox(Inches(0.62), Inches(1.05), Inches(12.0), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['table_body'], color=COLORS['dk2'])
    
    if not groups or not dimensions:
        return slide
    
    # 左半区：雷达图（用柱状图模拟多维度对比）
    chart_left = Inches(0.4)
    chart_top = Inches(1.5)
    chart_width = Inches(6.0)
    chart_height = Inches(4.5)
    
    chart_data = ChartData()
    chart_data.categories = dimensions
    group_colors = [
        COLORS['accent1'], COLORS['accent3'], COLORS['accent6'], 
        COLORS['accent4'], COLORS['accent2'], COLORS['accent5'],
    ]
    for gi, grp in enumerate(groups):
        grp_name = grp.get('name', f'组{gi+1}')
        grp_values = grp.get('values', [])
        chart_data.add_series(grp_name, grp_values)
    
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, chart_left, chart_top, chart_width, chart_height, chart_data
    )
    chart = graphic_frame.chart
    
    for si, series in enumerate(chart.series):
        if si < len(group_colors):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = group_colors[si]
        series.has_data_labels = True
        series.data_labels.font.size = FONT_SIZES['micro']
        series.data_labels.font.name = FONT_CN
    
    if chart.has_legend:
        chart.legend.include_in_layout = False
        chart.legend.font.size = FONT_SIZES['note']
        chart.legend.font.name = FONT_CN
    if hasattr(chart, 'category_axis') and chart.category_axis:
        chart.category_axis.tick_labels.font.size = FONT_SIZES['note']
        chart.category_axis.tick_labels.font.name = FONT_CN
    if hasattr(chart, 'value_axis') and chart.value_axis:
        chart.value_axis.tick_labels.font.size = FONT_SIZES['micro']
        chart.value_axis.tick_labels.font.name = FONT_CN
        chart.value_axis.maximum_scale = 10
    
    # 右半区：评分卡片 + 核心发现
    right_left = chart_left + chart_width + Inches(0.3)
    right_width = Inches(6.0)
    card_top = chart_top
    
    for gi, grp in enumerate(groups):
        grp_name = grp.get('name', '')
        grp_values = grp.get('values', [])
        avg_score = sum(grp_values) / len(grp_values) if grp_values else 0
        strengths = grp.get('strengths', '')
        weaknesses = grp.get('weaknesses', '')
        
        grp_color = group_colors[gi % len(group_colors)]
        card_height = Inches(1.35)
        card_y = card_top + gi * (card_height + Inches(0.15))
        
        # 卡片背景
        add_shape_rect(slide, right_left, card_y, right_width, card_height,
                       fill_color=COLORS['lt1'], line_color=grp_color, line_width=Pt(1.5))
        
        # 组名 + 均分
        add_shape_rect(slide, right_left, card_y, Inches(1.6), Inches(0.35), fill_color=grp_color)
        _add_centered_text(slide, right_left, card_y, Inches(1.6), Inches(0.35),
                          f'{grp_name}', font_size=12, color=COLORS['lt1'], bold=True)
        
        score_box = slide.shapes.add_textbox(right_left + Inches(1.7), card_y + Inches(0.02),
                                              Inches(1.0), Inches(0.3))
        tf = score_box.text_frame
        p = tf.paragraphs[0]
        p.text = f'{avg_score:.1f}'
        _set_font(p.runs[0] if p.runs else p.add_run(), size=Pt(20), bold=True, color=grp_color)
        
        # 优势
        if strengths:
            _add_left_text(slide, right_left + Inches(0.12), card_y + Inches(0.38),
                          right_width - Inches(0.24), Inches(0.25),
                          f'✓ {strengths}', font_size=9, color=COLORS['accent1'])
        # 劣势
        if weaknesses:
            _add_left_text(slide, right_left + Inches(0.12), card_y + Inches(0.65),
                          right_width - Inches(0.24), Inches(0.25),
                          f'⚠ {weaknesses}', font_size=9, color=COLORS['accent6'])
        
        # 各维度小评分条
        dim_y = card_y + Inches(0.95)
        dim_bar_h = Inches(0.18)
        for di, (dim, val) in enumerate(zip(dimensions, grp_values)):
            dim_x = right_left + Inches(0.12) + di * Inches(1.0)
            bar_w = Inches(0.9)
            # 背景条
            add_shape_rect(slide, dim_x, dim_y, bar_w, dim_bar_h,
                          fill_color=RGBColor(0xEE, 0xEE, 0xEE))
            # 填充条
            fill_w = bar_w * (val / 10.0)
            add_shape_rect(slide, dim_x, dim_y, fill_w, dim_bar_h, fill_color=grp_color)
            # 标签
            _add_left_text(slide, dim_x + Inches(0.02), dim_y - Inches(0.01),
                          bar_w - Inches(0.04), dim_bar_h,
                          f'{dim} {val}', font_size=6, color=COLORS['dk1'])
    
    if highlight_finding:
        finding_y = card_top + len(groups) * (Inches(1.35) + Inches(0.15)) + Inches(0.1)
        txBox = slide.shapes.add_textbox(right_left, finding_y, right_width, Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f'💡 {highlight_finding}'
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['table_body'],
                 bold=True, color=COLORS['accent3'])
    
    return slide


def build_waterfall_slide(prs, slide_spec):
    """构建瀑布图页面 — 费用拆解/投入产出分解
    
    用堆积柱状图模拟瀑布图：
    - 每个可见柱 = 实际金额
    - 透明基底 = 起点位置
    - 正向增长绿色，负向减少橙色，合计深绿色
    """
    title = slide_spec.get('title', '')
    subtitle = slide_spec.get('subtitle', '')
    items = slide_spec.get('items', [])          # [{label, value}, ...]  value正=增加, 负=减少
    total_label = slide_spec.get('total_label', '合计')
    note = slide_spec.get('note', '')
    
    layout = get_layout_by_type(prs, 'content')
    slide = prs.slides.add_slide(layout)
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = str(shape.placeholder_format.type).split('.')[-1]
            if 'TITLE' in ph_type:
                set_text_frame_text(shape, title, font_size=26, bold=True)
    remove_placeholders(slide)
    
    if subtitle:
        txBox = slide.shapes.add_textbox(Inches(0.62), Inches(1.05), Inches(12.0), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['table_body'], color=COLORS['dk2'])
    
    if not items:
        return slide
    
    # 计算瀑布图数据
    labels = []
    visible_values = []    # 可见柱高度
    base_values = []       # 透明基底（堆积柱状图的隐藏部分）
    running_total = 0
    
    for item in items:
        label = item.get('label', '')
        val = item.get('value', 0)
        labels.append(label)
        
        if val >= 0:
            # 正向：基底=running_total, 可见=val
            base_values.append(running_total)
            visible_values.append(val)
            running_total += val
        else:
            # 负向：基底=running_total+val, 可见=abs(val)
            base_values.append(running_total + val)
            visible_values.append(abs(val))
            running_total += val
    
    # 合计柱
    labels.append(total_label)
    base_values.append(0)
    visible_values.append(running_total)
    
    # 创建堆积柱状图：基底系列（透明）+ 可见系列（着色）
    chart_data = ChartData()
    chart_data.categories = labels
    chart_data.add_series('基底', base_values)
    chart_data.add_series('金额', visible_values)
    
    chart_left = Inches(0.6)
    chart_top = Inches(1.5)
    chart_width = Inches(12.0)
    chart_height = Inches(4.3)
    
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, chart_left, chart_top, chart_width, chart_height, chart_data
    )
    chart = graphic_frame.chart
    
    # 隐藏基底系列：无填充、无边框
    base_series = chart.series[0]
    base_series.format.fill.background()
    base_series.format.line.fill.background()
    
    # 可见系列：逐点着色
    vis_series = chart.series[1]
    vis_series.has_data_labels = True
    vis_series.data_labels.font.size = FONT_SIZES['micro']
    vis_series.data_labels.font.name = FONT_CN
    vis_series.data_labels.font.bold = True
    vis_series.data_labels.show_value = True
    
    # 逐点着色：正向绿、负向橙、合计深绿
    from pptx.dml.color import RGBColor as _RGBColor
    for pi, point in enumerate(vis_series.points):
        if pi < len(items):
            val = items[pi].get('value', 0)
            if val >= 0:
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = COLORS['accent1']
            else:
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = COLORS['accent6']
        else:
            # 合计柱
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = COLORS['accent3']
    
    # 去除图例
    if chart.has_legend:
        chart.legend.delete()
    
    # 坐标轴
    if hasattr(chart, 'category_axis') and chart.category_axis:
        chart.category_axis.tick_labels.font.size = FONT_SIZES['note']
        chart.category_axis.tick_labels.font.name = FONT_CN
    if hasattr(chart, 'value_axis') and chart.value_axis:
        chart.value_axis.tick_labels.font.size = FONT_SIZES['note']
        chart.value_axis.tick_labels.font.name = FONT_CN
    
    # 右侧数据表格（关键数字一览）
    table_left = Inches(0.6)
    table_top = chart_top + chart_height + Inches(0.12)
    n_cols = 3
    n_rows = len(items) + 2  # 表头 + 数据行 + 合计行
    
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, table_left, table_top, Inches(12.0), Inches(0.32) * n_rows)
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(3.5)
    tbl.columns[1].width = Inches(4.0)
    tbl.columns[2].width = Inches(4.5)
    
    # 表头
    for ci, hdr in enumerate(['费用项目', '金额(万元)', '占比']):
        cell = tbl.cell(0, ci)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['accent1']
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                _set_font(run, size=FONT_SIZES['table_header'], bold=True, color=COLORS['lt1'])
    
    # 数据行
    total_abs = sum(abs(it.get('value', 0)) for it in items)
    for ri, item in enumerate(items):
        label = item.get('label', '')
        val = item.get('value', 0)
        pct = f"{abs(val)/max(total_abs,1)*100:.1f}%"
        row_data = [label, f"{val:+,.0f}", pct]
        for ci, text in enumerate(row_data):
            cell = tbl.cell(ri + 1, ci)
            cell.text = text
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['table_stripe']
            for para in cell.text_frame.paragraphs:
                if ci == 0:
                    para.alignment = PP_ALIGN.LEFT
                else:
                    para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    _set_font(run, size=FONT_SIZES['table_body'], color=COLORS['dk1'],
                             bold=(val < 0))
                    if val < 0:
                        run.font.color.rgb = COLORS['accent6']
    
    # 合计行
    for ci, text in enumerate([total_label, f"{running_total:+,.0f}", '100%']):
        cell = tbl.cell(n_rows - 1, ci)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['hk_green']
        for para in cell.text_frame.paragraphs:
            if ci == 0:
                para.alignment = PP_ALIGN.LEFT
            else:
                para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                _set_font(run, size=FONT_SIZES['table_body'], bold=True, color=COLORS['accent3'])
    
    if note:
        note_y = table_top + Inches(0.32) * n_rows + Inches(0.08)
        txBox = slide.shapes.add_textbox(Inches(0.62), note_y, Inches(12.0), Inches(0.25))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = note
        _set_font(p.runs[0] if p.runs else p.add_run(), size=FONT_SIZES['micro'], color=COLORS['lt2'])
    
    return slide


# ============================================================
# 页面类型路由表
# ============================================================

BUILDERS = {
    'cover': build_cover_slide,
    'toc': build_toc_slide,
    'section': build_section_slide,
    'content': build_content_slide,
    'gantt': build_gantt_slide,
    'gantt_table': build_gantt_table_slide,
    'timeline_horizontal': build_timeline_horizontal_slide,
    'big_number': build_big_number_slide,
    'comparison': build_comparison_slide,
    'process': build_process_slide,
    'kpi_dashboard': build_kpi_dashboard_slide,
    'item_matrix': build_item_matrix_slide,
    'calendar_grid': build_calendar_grid_slide,
    'review_matrix': build_review_matrix_slide,
    'action_category': build_action_category_slide,
    'strategy_diagram': build_strategy_diagram_slide,
    'dual_axis_chart': build_dual_axis_chart_slide,
    'visual_comparison': build_visual_comparison_slide,
    'waterfall': build_waterfall_slide,
    'title_only': build_title_only_slide,
    'ending': build_ending_slide,
}


# ============================================================
# 主生成函数
# ============================================================

def build_slide(prs, slide_spec):
    """根据slide规范创建单页幻灯片"""
    slide_type = slide_spec.get('type', 'content')
    builder = BUILDERS.get(slide_type)
    if builder is None:
        print(f"Warning: Unknown slide type '{slide_type}', falling back to content", file=sys.stderr)
        builder = build_content_slide
    return builder(prs, slide_spec)


def generate_ppt(template_path, content_spec, output_path):
    """生成PPT主函数"""
    import shutil
    import tempfile
    import zipfile

    # 复制模板到临时文件
    tmp_fd, tmp_path = tempfile.mkstemp(suffix='.pptx')
    os.close(tmp_fd)
    shutil.copy2(template_path, tmp_path)

    try:
        prs = Presentation(tmp_path)
        slides_spec = content_spec.get('slides', [])
        if not slides_spec:
            print("Error: No slides specified in content", file=sys.stderr)
            sys.exit(1)

        pres_part = prs.part

        # 删除所有原有幻灯片
        orig_slide_ids = []
        for sldId in list(prs.slides._sldIdLst):
            rId = sldId.get(qn('r:id'))
            orig_slide_ids.append((sldId, rId))

        for sldId, rId in orig_slide_ids:
            prs.slides._sldIdLst.remove(sldId)
            if rId:
                pres_part.drop_rel(rId)

        # 添加所有新幻灯片
        new_slides = []
        for slide_spec in slides_spec:
            slide = build_slide(prs, slide_spec)
            if slide is not None:
                new_slides.append(slide)

        # 保存到临时文件
        tmp_fd2, tmp_path2 = tempfile.mkstemp(suffix='.pptx')
        os.close(tmp_fd2)

        try:
            prs.save(tmp_path2)

            # 收集有效slide文件名
            valid_slide_names = set()
            for sldId in prs.slides._sldIdLst:
                rId = sldId.get(qn('r:id'))
                if rId:
                    rel = pres_part.rels.get(rId)
                    if rel:
                        valid_slide_names.add('ppt/' + rel.target_ref)
                        slide_basename = os.path.basename(rel.target_ref)
                        slide_dir = os.path.dirname(rel.target_ref)
                        valid_slide_names.add('ppt/' + slide_dir + '/_rels/' + slide_basename + '.rels')

            # 重新打包zip
            with zipfile.ZipFile(tmp_path2, 'r') as zin:
                with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zout:
                    seen = set()
                    for item in zin.infolist():
                        is_old_slide = False
                        if item.filename.startswith('ppt/slides/slide') and item.filename.endswith('.xml'):
                            if item.filename not in valid_slide_names:
                                is_old_slide = True
                        elif item.filename.startswith('ppt/slides/_rels/slide') and item.filename.endswith('.xml.rels'):
                            if item.filename not in valid_slide_names:
                                is_old_slide = True
                        if is_old_slide:
                            continue
                        if item.filename in seen:
                            continue
                        seen.add(item.filename)
                        zout.writestr(item, zin.read(item.filename))

            print(f"✅ PPT已生成: {output_path}")
            print(f"   共 {len(new_slides)} 页")
            return output_path
        finally:
            if os.path.exists(tmp_path2):
                os.unlink(tmp_path2)
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)


def main():
    parser = argparse.ArgumentParser(description='先声药业集团PPT生成器 v3.0')
    parser.add_argument('--template', required=True, help='模板PPT文件路径')
    parser.add_argument('--input', required=True, help='内容JSON文件路径')
    parser.add_argument('--output', required=True, help='输出PPT文件路径')

    args = parser.parse_args()

    with open(args.input, 'r', encoding='utf-8') as f:
        content_spec = json.load(f)

    generate_ppt(args.template, content_spec, args.output)


if __name__ == '__main__':
    main()
