"""
先声药业 PPT 排版引擎
融合：官方母版模板 + 科唯可实战风格 + 品牌提示物矩阵风格 + 营销管理学术风格

精确参数来源于对以下文件的深度分析：
- 10-版本号202601集团PPT模板.pptx（官方母版）
- 科唯可-汇报任董0528.pptx（实战参考）
- 品牌提示物使用情况概览.pptx（矩阵布局参考）
- 何为专业的营销管理-任董.pptx（文字密集型参考）
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy
import os

# ============================================================
# 一、官方模板主题色（从 theme.xml 提取）
# ============================================================
COLORS = {
    # 主题色（官方模板）
    'accent1':      RGBColor(0x00, 0xB0, 0x52),  # 先声绿 - 主色
    'accent2':      RGBColor(0x8F, 0xD4, 0x00),  # 亮绿
    'accent3':      RGBColor(0x00, 0x66, 0x47),  # 深绿
    'accent4':      RGBColor(0x00, 0xB5, 0xBD),  # 青色
    'accent5':      RGBColor(0xFC, 0xC9, 0x17),  # 黄色
    'accent6':      RGBColor(0xF5, 0x66, 0x00),  # 橙色
    
    # 文字色
    'dk1':          RGBColor(0x33, 0x33, 0x33),  # 正文黑
    'dk2':          RGBColor(0x44, 0x54, 0x6A),  # 次级黑
    'lt1':          RGBColor(0xFF, 0xFF, 0xFF),  # 白色
    'lt2':          RGBColor(0xE7, 0xE6, 0xE6),  # 浅灰
    
    # 业务色（从参考PPT提取）
    'table_header': RGBColor(0x00, 0xB0, 0x52),  # 表头绿
    'table_stripe': RGBColor(0xF0, 0xF7, 0xF4),  # 隔行浅绿
    'bottom_bar':   RGBColor(0x00, 0xB0, 0x52),  # 底部色块
    'tag_green':    RGBColor(0x02, 0x9B, 0x46),  # 标签深绿（品牌提示物风格）
    'tag_light_green': RGBColor(0x5A, 0xC7, 0x1E),  # 标签亮绿
    'score_orange': RGBColor(0xE4, 0x80, 0x30),  # 评分橙
    'new_item':     RGBColor(0x80, 0x80, 0x80),  # 灰色（表示待替换/未激活项）
    
    # 背景色（用于行/卡片背景）
    'hk_green':     RGBColor(0xE8, 0xF5, 0xEE),  # 浅绿背景
    'hlk_green':    RGBColor(0xC8, 0xEB, 0xD6),  # 高亮浅绿
    'hlk_orange':   RGBColor(0xFF, 0xED, 0xE1),  # 浅橙背景
    'orange':       RGBColor(0xE4, 0x80, 0x30),  # 橙色（别名）
    
    # 架构图配色（科唯可参考）
    'strategy_dark':  RGBColor(0x0A, 0xA4, 0x47),  # 深绿
    'strategy_mid':   RGBColor(0x38, 0xB7, 0x2F),  # 中绿
    'strategy_light': RGBColor(0x69, 0xCA, 0x16),  # 浅绿
    'strategy_lighter': RGBColor(0x9C, 0xD3, 0x03), # 极浅绿
}

# 官方主题色名称映射（用于占位符继承主题色）
THEME_COLORS = {
    'ACCENT_1': COLORS['accent1'],
    'ACCENT_2': COLORS['accent2'],
    'ACCENT_3': COLORS['accent3'],
    'ACCENT_4': COLORS['accent4'],
    'ACCENT_5': COLORS['accent5'],
    'ACCENT_6': COLORS['accent6'],
    'TEXT_1': COLORS['dk1'],
    'BACKGROUND_1': COLORS['lt1'],
}

# ============================================================
# 二、字体规范（官方模板说明页确认）
# ============================================================
FONT_CN = '微软雅黑'
FONT_EN = 'Arial'

# 字号层级（融合官方模板 + 三份参考PPT）
FONT_SIZES = {
    # 封面
    'cover_title':      Pt(44),    # 封面大标题（营销管理参考）
    'cover_subtitle':   Pt(20),    # 封面副标题
    
    # 章节
    'section_title':    Pt(32),    # 节标题
    'section_subtitle': Pt(16),    # 节副标题
    
    # 内容页
    'page_tag':         Pt(20),    # 页面标签（如"复盘""反思与改进"）科唯可风格
    'page_title':       Pt(26),    # 页面标题（品牌提示物用 +mn-ea 26pt）
    'page_title_alt':   Pt(18),    # 页面标题备用
    
    # 正文层级
    'body_main':        Pt(16),    # 主正文（科唯可16pt）
    'body_normal':      Pt(14),    # 标准正文（科唯可14pt）
    'body_small':       Pt(12),    # 小正文（品牌提示物矩阵文字12pt）
    
    # 表格
    'table_header':     Pt(12),    # 表头
    'table_body':       Pt(11),    # 表格正文
    
    # 注释
    'note':             Pt(10),    # 注释
    'micro':            Pt(9),     # 微小注释（品牌提示物场景说明9pt）
    'tiny':             Pt(7),     # 极小注释（科唯可数据标注7pt）
    
    # 大数字
    'big_number':       Pt(36),    # KPI大数字
    'big_number_label': Pt(12),    # 大数字标签
    
    # 目录
    'toc_item':         Pt(18),    # 目录条目
}

# ============================================================
# 三、页面布局坐标（从官方模板母版提取，精确到英寸）
# ============================================================

# 幻灯片尺寸
SLIDE_WIDTH  = 13.333  # inches (12192000 EMU)
SLIDE_HEIGHT = 7.5     # inches (6858000 EMU)

# 封面页布局（官方模板 版式0）
COVER_LAYOUT = {
    'title':    (0.59, 1.05, 8.41, 0.79),     # CENTER_TITLE ph#0
    'subtitle': (0.63, 1.87, 8.37, 0.49),     # SUBTITLE ph#1
    'dept':     (0.63, 5.77, 3.94, 0.32),     # BODY ph#10
    'date':     (0.63, 6.14, 3.94, 0.32),     # BODY ph#11
}

# 目录页布局（官方模板 版式1 / 参考 科唯可 Slide2）
TOC_LAYOUT = {
    'content':  (3.84, 1.95, 8.76, 4.38),     # 目录文字区域
}

# 节标题页布局（官方模板 版式2）
SECTION_LAYOUT = {
    'title':    (1.63, 1.26, 8.33, 0.98),     # TITLE ph#0
    'body':     (0.73, 3.06, 6.04, 1.11),     # BODY ph#1
    'accent':   (0.73, 1.26, 0.61, 0.98),     # BODY ph#13 (左侧绿色竖条)
}

# 标准内容页布局（官方模板 版式3）—— 融合科唯可实战风格
CONTENT_LAYOUT = {
    'tag':      (0.62, 0.34, 9.89, 0.40),     # BODY ph#15 - 页面标签 20pt
    'title':    (0.62, 0.66, 9.89, 0.51),     # TITLE ph#0 - 页面标题
    'body':     (0.62, 1.35, 12.30, 5.30),    # 正文区域（含OBJECT占位符位置）
    'footer':   (0.63, 6.99, 4.53, 0.23),     # FOOTER ph#3
    'page_num': (12.50, 7.05, 0.56, 0.23),    # SLIDE_NUMBER ph#4
    
    # 底部色块区域（科唯可风格）
    'bottom_bar': (0.00, 5.90, 13.33, 1.60),  # 底部结论/总结区
    
    # 分栏布局
    'left_col':  (0.62, 1.35, 5.90, 5.30),    # 左栏
    'right_col': (6.80, 1.35, 6.00, 5.30),    # 右栏
}

# 仅标题页（官方模板 版式4）
TITLE_ONLY_LAYOUT = {
    'tag':      (0.62, 0.34, 9.89, 0.40),
    'title':    (0.62, 0.66, 9.89, 0.51),
    'body':     (0.62, 1.40, 12.30, 5.50),
    'page_num': (12.50, 7.05, 0.56, 0.23),
}

# 结束页（官方模板 版式7）
ENDING_LAYOUT = {
    'title':    (0.74, 1.90, 6.93, 0.98),     # CENTER_TITLE ph#0
    'body1':    (0.74, 4.73, 5.41, 0.32),     # BODY ph#10
    'body2':    (0.74, 5.05, 5.41, 0.34),     # BODY ph#18
}

# 矩阵布局（品牌提示物风格）
MATRIX_LAYOUT = {
    'scene_col':  (0.30, 1.50, 1.50, 5.80),   # 场景列
    'object_col': (1.80, 1.50, 2.30, 5.80),   # 对象列
    'items_col':  (4.10, 1.50, 5.50, 5.80),   # 物品列
    'cost_col':   (9.60, 1.50, 1.50, 5.80),   # 成本列
    'score_col':  (11.10, 1.50, 1.80, 5.80),  # 评分列
}

# 表格布局
TABLE_LAYOUT = {
    'full':      (0.52, 1.40, 12.46, 5.50),   # 全宽表格
    'with_text': (0.52, 3.00, 12.46, 4.00),   # 上方有文字的表格
    'comparison':(0.52, 1.40, 6.00, 5.50),    # 左侧对比表
}

# KPI仪表盘布局
KPI_LAYOUT = {
    'cards_start': (0.62, 1.60),              # 卡片起始位置
    'card_size':   (2.80, 1.50),              # 每个卡片尺寸
    'card_gap':    0.30,                       # 卡片间距
    'cards_per_row': 4,                        # 每行卡片数
}

# ============================================================
# 四、段落间距（融合四份PPT的排版参数）
# ============================================================
PARA_SPACING = {
    'title_before':     Pt(0),
    'title_after':      Pt(8),
    'body_before':      Pt(6),
    'body_after':       Pt(3),
    'compact_before':   Pt(2),
    'compact_after':    Pt(1),
    'section_before':   Pt(12),   # 段落间大间距
    'section_after':    Pt(6),
}

LINE_SPACING = {
    'title':    1.0,    # 标题单倍行距
    'body':     1.3,    # 正文1.3倍行距
    'compact':  1.15,   # 紧凑行距
    'table':    1.0,    # 表格单倍行距
}

# ============================================================
# 五、布局辅助函数
# ============================================================

def inches_to_emu(*args):
    """将英寸值转换为EMU"""
    if len(args) == 1:
        return int(args[0] * 914400)
    return tuple(int(a * 914400) for a in args)

def set_font(run, name_cn=FONT_CN, name_en=FONT_EN, size=None, bold=False, color=None):
    """统一设置字体"""
    run.font.name = name_cn
    # 设置东亚字体
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn('a:ea'), {})
    ea.set('typeface', name_cn)
    rPr.append(ea)
    # 设置拉丁字体
    latin = rPr.makeelement(qn('a:latin'), {})
    latin.set('typeface', name_en)
    rPr.append(latin)
    
    if size:
        run.font.size = size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color

def set_paragraph_spacing(para, before=None, after=None, line_spacing=None):
    """设置段落间距"""
    from lxml import etree as _etree
    pPr = para._pPr
    if pPr is None:
        pPr = _etree.SubElement(para._p, qn('a:pPr'))
    
    if before is not None:
        spcBef = pPr.find(qn('a:spcBef'))
        if spcBef is None:
            spcBef = _etree.SubElement(pPr, qn('a:spcBef'))
        spcPts = spcBef.find(qn('a:spcPts'))
        if spcPts is None:
            spcPts = _etree.SubElement(spcBef, qn('a:spcPts'))
        spcPts.set('val', str(int(before / 12700)))  # EMU to points
    
    if after is not None:
        spcAft = pPr.find(qn('a:spcAft'))
        if spcAft is None:
            spcAft = _etree.SubElement(pPr, qn('a:spcAft'))
        spcPts = spcAft.find(qn('a:spcPts'))
        if spcPts is None:
            spcPts = _etree.SubElement(spcAft, qn('a:spcPts'))
        spcPts.set('val', str(int(after / 12700)))
    
    if line_spacing is not None:
        lnSpc = pPr.find(qn('a:lnSpc'))
        if lnSpc is None:
            lnSpc = _etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = lnSpc.find(qn('a:spcPct'))
        if spcPct is None:
            spcPct = _etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 100000)))  # percentage in 1/1000%

def add_textbox(slide, left, top, width, height):
    """添加文本框"""
    from pptx.util import Inches
    return slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )

def add_paragraph(tf, text, font_size=None, bold=False, color=None, 
                  alignment=None, level=0, before=None, after=None, line_spacing=None):
    """添加段落并设置格式"""
    para = tf.add_paragraph()
    para.level = level
    run = para.add_run()
    run.text = text
    set_font(run, size=font_size, bold=bold, color=color)
    set_paragraph_spacing(para, before=before, after=after, line_spacing=line_spacing)
    if alignment is not None:
        para.alignment = alignment
    return para

def add_first_paragraph(tf, text, font_size=None, bold=False, color=None,
                        alignment=None, level=0, before=None, after=None, line_spacing=None):
    """替换文本框第一个段落"""
    para = tf.paragraphs[0]
    para.level = level
    run = para.add_run()
    run.text = text
    set_font(run, size=font_size, bold=bold, color=color)
    set_paragraph_spacing(para, before=before, after=after, line_spacing=line_spacing)
    if alignment is not None:
        para.alignment = alignment
    return para

def fill_placeholder(slide, ph_idx, text, font_size=None, bold=False, color=None):
    """填充占位符文本"""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == ph_idx:
            tf = shape.text_frame
            # 清除现有文本
            tf.clear()
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = text
            set_font(run, size=font_size, bold=bold, color=color)
            return shape
    return None

def get_placeholder(slide, ph_idx):
    """获取占位符shape"""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == ph_idx:
            return shape
    return None

def remove_placeholders(slide, types_to_remove=None):
    """删除指定类型的占位符（先收集再删除，避免遍历跳过）"""
    if types_to_remove is None:
        types_to_remove = ['BODY', 'OBJECT']
    
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = str(shape.placeholder_format.type)
            for t in types_to_remove:
                if t in ph_type:
                    shapes_to_remove.append(shape)
                    break
    
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

def remove_shape_safe(shape):
    """安全删除shape"""
    sp = shape._element
    sp.getparent().remove(sp)

def add_bottom_bar(slide, height=1.30, color=None):
    """添加底部色块（科唯可风格）"""
    if color is None:
        color = COLORS['bottom_bar']
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(SLIDE_HEIGHT - height),
        Inches(SLIDE_WIDTH), Inches(height)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()  # 无边框
    return bar

def add_table(slide, rows, cols, left, top, width, height):
    """添加真实表格对象"""
    tbl_shape = slide.shapes.add_table(
        rows, cols,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    return tbl_shape.table

def style_table_header(table, col_count):
    """设置表头样式：先声绿背景 + 白色粗体"""
    for ci in range(col_count):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['table_header']
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                set_font(run, size=FONT_SIZES['table_header'], bold=True, color=COLORS['lt1'])

def style_table_body(table, start_row, col_count):
    """设置表体样式：隔行浅绿"""
    for ri in range(start_row, len(table.rows)):
        for ci in range(col_count):
            cell = table.cell(ri, ci)
            if (ri - start_row) % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['table_stripe']
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    set_font(run, size=FONT_SIZES['table_body'], color=COLORS['dk1'])

def build_body_content(tf, items, default_size=None, default_color=None):
    """
    构建正文内容，支持：
    - 纯文本列表
    - 小标题+正文（BOLD标记）
    - 多级缩进
    """
    if default_size is None:
        default_size = FONT_SIZES['body_normal']
    if default_color is None:
        default_color = COLORS['dk1']
    
    for item in items:
        if isinstance(item, str):
            # 简单文本
            if item.startswith('## '):
                # 小标题
                add_paragraph(tf, item[3:], 
                            font_size=FONT_SIZES['body_main'], bold=True,
                            color=COLORS['accent1'],
                            before=PARA_SPACING['section_before'],
                            after=PARA_SPACING['section_after'])
            elif item.startswith('- '):
                # 列表项
                add_paragraph(tf, item, 
                            font_size=default_size, color=default_color,
                            before=PARA_SPACING['body_before'],
                            after=PARA_SPACING['body_after'],
                            line_spacing=LINE_SPACING['body'])
            else:
                add_paragraph(tf, item,
                            font_size=default_size, color=default_color,
                            before=PARA_SPACING['body_before'],
                            after=PARA_SPACING['body_after'],
                            line_spacing=LINE_SPACING['body'])
        elif isinstance(item, dict):
            # 结构化条目
            text = item.get('text', '')
            size = item.get('size', default_size)
            bold = item.get('bold', False)
            color = item.get('color', default_color)
            level = item.get('level', 0)
            before = item.get('before', PARA_SPACING['body_before'])
            after = item.get('after', PARA_SPACING['body_after'])
            
            add_paragraph(tf, text, font_size=size, bold=bold, color=color,
                         level=level, before=before, after=after,
                         line_spacing=LINE_SPACING['body'])

def add_kpi_card(slide, left, top, width, height, 
                 value, label, change=None, color=None):
    """添加KPI卡片"""
    if color is None:
        color = COLORS['accent1']
    
    # 卡片背景
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xF5, 0xFB, 0xF7)
    card.line.color.rgb = RGBColor(0xE0, 0xF0, 0xE5)
    
    tf = card.text_frame
    tf.word_wrap = True
    
    # 大数字
    add_first_paragraph(tf, str(value),
                       font_size=FONT_SIZES['big_number'], bold=True,
                       color=color, alignment=PP_ALIGN.CENTER)
    
    # 标签
    add_paragraph(tf, label,
                 font_size=FONT_SIZES['big_number_label'], color=COLORS['dk1'],
                 alignment=PP_ALIGN.CENTER)
    
    # 变化率
    if change:
        change_color = COLORS['accent1'] if change.startswith('+') else COLORS['accent6']
        add_paragraph(tf, change,
                     font_size=FONT_SIZES['note'], bold=True,
                     color=change_color, alignment=PP_ALIGN.CENTER)
    
    return card

def add_green_tag(slide, left, top, width, height, text, font_size=None, color=None):
    """添加绿色标签色块（品牌提示物风格）"""
    if color is None:
        color = COLORS['accent1']
    if font_size is None:
        font_size = FONT_SIZES['body_small']
    
    tag = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = color
    tag.line.fill.background()
    
    tf = tag.text_frame
    tf.word_wrap = True
    add_first_paragraph(tf, text, font_size=font_size, bold=True,
                       color=COLORS['lt1'], alignment=PP_ALIGN.CENTER)
    return tag

# ============================================================
# 六、页面类型布局索引
# ============================================================
LAYOUT_INDEX = {
    'cover':        0,   # 封面
    'toc':          1,   # 1_仅logo页（目录）
    'section':      2,   # 节标题
    'content':      3,   # 标题和内容
    'title_only':   4,   # 仅标题页
    'logo_only':    5,   # 仅logo页
    'blank':        6,   # 空白
    'ending':       7,   # 末尾幻灯片
}

def get_layout_name(page_type):
    """页面类型到版式名称映射"""
    mapping = {
        'cover': '封面',
        'toc': '1_仅logo页',
        'section': '节标题',
        'content': '标题和内容',
        'title_only': '仅标题页',
        'logo_only': '仅logo页',
        'blank': '空白',
        'ending': '末尾幻灯片',
    }
    return mapping.get(page_type, '标题和内容')

print("✅ simcere_layout.py 排版引擎加载完成")
print(f"   主题色: {len(COLORS)} 种")
print(f"   字号层级: {len(FONT_SIZES)} 级")
print(f"   布局模板: 7 种页面类型")
print(f"   字体: {FONT_CN} / {FONT_EN}")
